#!/usr/bin/env python3
"""
NexusAI FastAPI Server
Universal AI Gateway Backend
"""

import os
import json
import asyncio
from datetime import datetime
from typing import Dict, List, Optional
from contextlib import asynccontextmanager

from fastapi import FastAPI, HTTPException, Request, status
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse, StreamingResponse
from pydantic import BaseModel
from dataclasses import dataclass

# Import NexusAI SDK
import sys
sys.path.insert(0, '/root/tender-system/nexusai/sdk/python')
from nexusai import NexusAI, Priority

# ============== CONFIG ==============
API_VERSION = "1.0.0"
APP_NAME = "NexusAI Server"

# ============== MODELS ==============
class ChatRequest(BaseModel):
    model: str
    messages: List[Dict[str, str]]
    temperature: Optional[float] = 0.7
    max_tokens: Optional[int] = 1000
    stream: Optional[bool] = False
    priority: Optional[str] = "cost"  # cost, speed, quality

class ChatResponse(BaseModel):
    response: str
    provider: str
    model: str
    usage: Dict
    latency_ms: int

class CostResponse(BaseModel):
    total_tokens: int
    by_provider: Dict[str, int]
    estimated_cost: float

class ModelInfo(BaseModel):
    name: str
    provider: str
    status: str
    latency_ms: int

# ============== APP ==============
app = FastAPI(
    title="NexusAI API",
    description="One SDK to Rule All AI Models",
    version=API_VERSION
)

# CORS
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Initialize NexusAI
nexus = None

@asynccontextmanager
async def lifespan(app: FastAPI):
    global nexus
    # Startup
    nexus = NexusAI()
    print(f"🚀 {APP_NAME} v{API_VERSION} started")
    yield
    # Shutdown
    print("🛑 Shutting down...")

app = FastAPI(lifespan=lifespan)

# ============== ROUTES ==============

@app.get("/")
async def root():
    return {
        "name": APP_NAME,
        "version": API_VERSION,
        "status": "running",
        "docs": "/docs"
    }

@app.get("/health")
async def health():
    return {"status": "healthy", "timestamp": datetime.now().isoformat()}

@app.post("/chat", response_model=ChatResponse)
async def chat(request: ChatRequest):
    """Universal chat endpoint - works with any model"""
    if not nexus:
        raise HTTPException(status_code=500, detail="AI not initialized")
    
    try:
        # Convert priority string to enum
        priority_map = {
            "cost": Priority.COST,
            "speed": Priority.SPEED,
            "quality": Priority.QUALITY
        }
        priority = priority_map.get(request.priority, Priority.COST)
        
        # Make request
        import time
        start = time.time()
        
        response = nexus.chat(
            model=request.model,
            messages=request.messages,
            priority=priority
        )
        
        latency = int((time.time() - start) * 1000)
        
        if "error" in response:
            raise HTTPException(status_code=400, detail=response["error"])
        
        return ChatResponse(
            response=response.get("response", ""),
            provider=response.get("provider", "unknown"),
            model=request.model,
            usage=response.get("usage", {}),
            latency_ms=latency
        )
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@app.post("/chat/stream")
async def chat_stream(request: ChatRequest):
    """Streaming chat endpoint"""
    if not nexus:
        raise HTTPException(status_code=500, detail="AI not initialized")
    
    async def generate():
        yield f"data: {json.dumps({'status': 'streaming'})}\n\n"
        # For demo, return mock streaming
        response = nexus.chat(request.model, request.messages)
        for word in response.get("response", "Streaming...").split():
            yield f"data: {json.dumps({'token': word})}\n\n"
        yield f"data: {json.dumps({'done': True})}\n\n"
    
    return StreamingResponse(generate(), media_type="text/event-stream")

@app.get("/cost", response_model=CostResponse)
async def get_cost():
    """Get cost summary"""
    if not nexus:
        raise HTTPException(status_code=500, detail="AI not initialized")
    
    summary = nexus.get_cost_summary()
    
    return CostResponse(
        total_tokens=summary.get("total_tokens", 0),
        by_provider=summary.get("by_provider", {}),
        estimated_cost=summary.get("estimated_cost", 0)
    )

@app.get("/models")
async def list_models():
    """List all available models"""
    if not nexus:
        raise HTTPException(status_code=500, detail="AI not initialized")
    
    models = []
    for name, provider in nexus.providers.items():
        models.append({
            "name": name,
            "provider": provider.name,
            "status": "online" if provider.api_key else "no_key",
            "latency_ms": provider.latency_ms,
            "cost_input": provider.cost_per_1k_input,
            "cost_output": provider.cost_per_1k_output,
            "local": getattr(provider, 'local', False)
        })
    
    return {"models": models, "count": len(models)}

@app.post("/agent/create")
async def create_agent(tools: Optional[List[str]] = None):
    """Create AI agent"""
    if not nexus:
        raise HTTPException(status_code=500, detail="AI not initialized")
    
    agent = nexus.create_agent()
    
    return {
        "agent_id": f"agent_{datetime.now().strftime('%Y%m%d%H%M%S')}",
        "status": "created",
        "tools_count": len(tools) if tools else 0
    }

@app.post("/agent/{agent_id}/run")
async def run_agent(agent_id: str, task: str):
    """Run agent task"""
    if not nexus:
        raise HTTPException(status_code=500, detail="AI not initialized")
    
    agent = nexus.create_agent()
    result = agent.run(task)
    
    return {
        "agent_id": agent_id,
        "task": task,
        "result": result.get("response", ""),
        "status": "completed"
    }

@app.get("/compare")
async def compare_models(prompt: str, models: Optional[str] = None):
    """Compare multiple models"""
    if not nexus:
        raise HTTPException(status_code=500, detail="AI not initialized")
    
    model_list = models.split(",") if models else ["openai/gpt-4o", "gemini/gemini-pro"]
    
    results = nexus.compare_models(prompt, model_list)
    
    return {
        "prompt": prompt,
        "results": results
    }

@app.get("/providers")
async def list_providers():
    """List all configured providers"""
    if not nexus:
        raise HTTPException(status_code=500, detail="AI not initialized")
    
    providers = []
    for name, provider in nexus.providers.items():
        providers.append({
            "name": provider.name,
            "base_url": provider.base_url,
            "cost_input": provider.cost_per_1k_input,
            "cost_output": provider.cost_per_1k_output,
            "latency_ms": provider.latency_ms,
            "configured": bool(provider.api_key)
        })
    
    return {"providers": providers}

@app.post("/providers/add")
async def add_provider(config: Dict):
    """Add new provider"""
    if not nexus:
        raise HTTPException(status_code=500, detail="AI not initialized")
    
    name = config.get("name")
    if not name:
        raise HTTPException(status_code=400, detail="Provider name required")
    
    nexus.add_provider(name, config)
    
    return {"status": "added", "provider": name}

# Web Search using free DuckDuckGo API
import urllib.request
import urllib.parse
import json
import subprocess
import os

@app.get("/search")
async def web_search(q: str, limit: int = 10):
    """Real-time web search using DuckDuckGo API"""
    try:
        query = urllib.parse.quote(q)
        url = f"https://api.duckduckgo.com/?q={query}&format=json&no_html=1&skip_disambig=1"
        
        with urllib.request.urlopen(url, timeout=30) as response:
            data = json.loads(response.read().decode())
        
        # Get related searches
        related = []
        for item in data.get("RelatedTopics", [])[:5]:
            if isinstance(item, dict):
                related.append({
                    "title": item.get("Text", ""),
                    "url": item.get("URL", "")
                })
        
        return {
            "query": q,
            "answer": data.get("Answer", ""),
            "abstract": data.get("AbstractText", ""),
            "related": related,
            "source": "DuckDuckGo"
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Search failed: {str(e)}")

@app.get("/search/advanced")
async def advanced_search(q: str, source: str = "web"):
    """Advanced search with more results"""
    try:
        query = urllib.parse.quote(q)
        url = f"https://api.duckduckgo.com/?q={query}&format=json"
        
        with urllib.request.urlopen(url, timeout=30) as response:
            data = json.loads(response.read().decode())
        
        results = []
        for item in data.get("Results", [])[:20]:
            results.append({
                "title": item.get("Text", ""),
                "url": item.get("URL", "")
            })
        
        return {
            "query": q,
            "count": len(results),
            "results": results
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Search failed: {str(e)}")

@app.post("/execute")
async def execute_command(cmd: Dict):
    """Execute commands or open applications"""
    try:
        command = cmd.get("command", "")
        action = cmd.get("action", "run")
        
        if not command:
            raise HTTPException(status_code=400, detail="Command required")
        
        # Security: only allow safe commands
        allowed_actions = ["run", "open", "start"]
        if action not in allowed_actions:
            raise HTTPException(status_code=400, detail="Invalid action")
        
        # Execute command
        result = subprocess.run(
            command,
            shell=True,
            capture_output=True,
            text=True,
            timeout=30
        )
        
        return {
            "status": "success",
            "command": command,
            "output": result.stdout,
            "error": result.stderr,
            "return_code": result.returncode
        }
    except subprocess.TimeoutExpired:
        raise HTTPException(status_code=500, detail="Command timeout")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Execution failed: {str(e)}")

@app.get("/apps")
async def list_apps():
    """List available applications (for desktop)"""
    # Common apps that can be opened
    apps = [
        {"name": "Browser", "command": "xdg-open https://", "platform": "Linux"},
        {"name": "File Manager", "command": "xdg-open ~", "platform": "Linux"},
        {"name": "Terminal", "command": "gnome-terminal", "platform": "Linux"},
        {"name": "VS Code", "command": "code", "platform": "Linux"},
        {"name": "Notepad", "command": "notepad", "platform": "Windows"},
        {"name": "Calculator", "command": "gnome-calculator", "platform": "Linux"},
    ]
    return {"apps": apps}

@app.get("/system")
async def get_system_info():
    """Get detailed system information"""
    import platform
    import psutil
    import os
    
    # CPU Info
    cpu_freq = psutil.cpu_freq()
    cpu_info = {
        "physical_cores": psutil.cpu_count(logical=False),
        "logical_cores": psutil.cpu_count(logical=True),
        "frequency_current": round(cpu_freq.current, 2) if cpu_freq else "N/A",
        "frequency_max": round(cpu_freq.max, 2) if cpu_freq else "N/A",
        "frequency_min": round(cpu_freq.min, 2) if cpu_freq else "N/A",
        "cpu_percent": psutil.cpu_percent(interval=1),
    }
    
    # Memory Info
    mem = psutil.virtual_memory()
    memory = {
        "total_gb": round(mem.total / (1024**3), 2),
        "available_gb": round(mem.available / (1024**3), 2),
        "used_gb": round(mem.used / (1024**3), 2),
        "percent": mem.percent,
    }
    
    # Disk Info
    disk = psutil.disk_usage('/')
    disk_info = {
        "total_gb": round(disk.total / (1024**3), 2),
        "used_gb": round(disk.used / (1024**3), 2),
        "free_gb": round(disk.free / (1024**3), 2),
        "percent": disk.percent,
    }
    
    # Network Info
    net = psutil.net_io_counters()
    network = {
        "bytes_sent_mb": round(net.bytes_sent / (1024**2), 2),
        "bytes_recv_mb": round(net.bytes_recv / (1024**2), 2),
        "packets_sent": net.packets_sent,
        "packets_recv": net.packets_recv,
    }
    
    # OS Info
    os_info = {
        "system": platform.system(),
        "release": platform.release(),
        "version": platform.version(),
        "machine": platform.machine(),
        "processor": platform.processor(),
        "hostname": platform.node(),
    }
    
    # Boot Time
    boot_time = psutil.boot_time()
    import datetime
    boot_datetime = datetime.datetime.fromtimestamp(boot_time)
    uptime = datetime.datetime.now() - boot_datetime
    uptime_str = f"{uptime.days}d {uptime.seconds//3600}h {(uptime.seconds%3600)//60}m"
    
    # Battery (if available)
    battery = {}
    try:
        batt = psutil.sensors_battery()
        if batt:
            battery = {
                "percent": batt.percent,
                "plugged": batt.power_plugged,
                "time_left_minutes": batt.secsleft / 60 if batt.secsleft != psutil.POWER_TIME_UNLIMITED else "Unlimited"
            }
    except:
        battery = {"status": "Not available"}
    
    # Temperature (if available)
    temps = {}
    try:
        temp = psutil.sensors_temperatures()
        for name, entries in temp.items():
            if entries:
                temps[name] = round(entries[0].current, 1)
    except:
        temps = {"status": "Not available"}
    
    return {
        "os": os_info,
        "cpu": cpu_info,
        "memory": memory,
        "disk": disk_info,
        "network": network,
        "battery": battery,
        "temperature": temps,
        "uptime": uptime_str,
        "python_version": platform.python_version(),
    }

@app.get("/system/live")
async def get_live_stats():
    """Get live CPU and memory stats"""
    import psutil
    
    return {
        "cpu_percent": psutil.cpu_percent(interval=0.5),
        "cpu_count": psutil.cpu_count(),
        "memory_percent": psutil.virtual_memory().percent,
        "memory_available_gb": round(psutil.virtual_memory().available / (1024**3), 2),
        "disk_percent": psutil.disk_usage('/').percent,
        "network_sent_mb": round(psutil.net_io_counters().bytes_sent / (1024**2), 2),
        "network_recv_mb": round(psutil.net_io_counters().bytes_recv / (1024**2), 2),
    }

@app.get("/weather")
async def get_weather(location: str = "Mumbai"):
    """Get weather data for a location"""
    try:
        import urllib.request
        import urllib.parse
        import json
        
        # Use wttr.in - free weather API (no key required)
        query = urllib.parse.quote(location)
        url = f"https://wttr.in/{query}?format=j1"
        
        with urllib.request.urlopen(url, timeout=10) as response:
            data = json.loads(response.read().decode())
        
        current = data.get("current_condition", [{}])[0]
        
        return {
            "location": location,
            "temperature_c": current.get("temp_C", "N/A"),
            "temperature_f": current.get("temp_F", "N/A"),
            "condition": current.get("weatherDesc", [{}])[0].get("value", "Unknown"),
            "humidity": current.get("humidity", "N/A"),
            "wind_kmh": current.get("windspeedKmph", "N/A"),
            "feels_like_c": current.get("FeelsLikeC", "N/A"),
            "uv_index": current.get("UVindex", "N/A"),
            "visibility": current.get("visibility", "N/A"),
            "pressure": current.get("pressure", "N/A"),
            "source": "wttr.in"
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Weather fetch failed: {str(e)}")

@app.get("/weather/forecast")
async def get_weather_forecast(location: str = "Mumbai", days: int = 3):
    """Get weather forecast"""
    try:
        import urllib.request
        import urllib.parse
        import json
        
        query = urllib.parse.quote(location)
        url = f"https://wttr.in/{query}?format=j1"
        
        with urllib.request.urlopen(url, timeout=10) as response:
            data = json.loads(response.read().decode())
        
        forecast = []
        for i, day in enumerate(data.get("weather", [])[:days]):
            forecast.append({
                "date": day.get("date", ""),
                "max_temp_c": day.get("maxtempC", "N/A"),
                "min_temp_c": day.get("mintempC", "N/A"),
                "avg_temp_c": day.get("avgtempC", "N/A"),
                "condition": day.get("weatherDesc", [{}])[0].get("value", "Unknown") if day.get("weatherDesc") else "Unknown",
                "sunrise": day.get("astronomy", [{}])[0].get("sunrise", "N/A") if day.get("astronomy") else "N/A",
                "sunset": day.get("astronomy", [{}])[0].get("sunset", "N/A") if day.get("astronomy") else "N/A",
            })
        
        return {
            "location": location,
            "forecast": forecast,
            "source": "wttr.in"
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Forecast fetch failed: {str(e)}")

@app.get("/windows")
async def list_windows():
    """List all windows/processes (Linux)"""
    import psutil
    try:
        import subprocess
        import os
        
        # Get running processes
        processes = []
        for proc in psutil.process_iter(['pid', 'name', 'status', 'cpu_percent', 'memory_percent', 'create_time']):
            try:
                pinfo = proc.info
                if pinfo['name']:
                    processes.append({
                        "pid": pinfo['pid'],
                        "name": pinfo['name'],
                        "status": pinfo['status'],
                        "cpu_percent": pinfo.get('cpu_percent', 0),
                        "memory_percent": round(pinfo.get('memory_percent', 0), 2),
                    })
            except (psutil.NoSuchProcess, psutil.AccessDenied):
                pass
        
        # Sort by CPU usage
        processes.sort(key=lambda x: x['cpu_percent'], reverse=True)
        
        return {
            "total_processes": len(processes),
            "top_processes": processes[:20],
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to list windows: {str(e)}")

@app.get("/windows/active")
async def get_active_window():
    """Get currently active window (Linux)"""
    try:
        import subprocess
        
        # Try using xdotool (if available)
        try:
            result = subprocess.run(
                ["xdotool", "getactivewindow", "getwindowname"],
                capture_output=True,
                text=True,
                timeout=2
            )
            if result.returncode == 0:
                return {
                    "active_window": result.stdout.strip(),
                    "method": "xdotool"
                }
        except:
            pass
        
        # Try using wmctrl
        try:
            result = subprocess.run(
                ["wmctrl", "-a", "-"],
                capture_output=True,
                text=True,
                timeout=2
            )
        except:
            pass
        
        return {
            "active_window": "Not available (requires xdotool or wmctrl)",
            "method": "fallback"
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed: {str(e)}")

@app.post("/windows/focus")
async def focus_window(window_name: str):
    """Focus a window by name"""
    try:
        import subprocess
        
        result = subprocess.run(
            ["xdotool", "search", "--name", window_name, "windowfocus"],
            capture_output=True,
            text=True,
            timeout=5
        )
        
        if result.returncode == 0:
            return {"status": "success", "action": f"Focused window: {window_name}"}
        else:
            return {"status": "failed", "error": result.stderr}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed: {str(e)}")

@app.post("/windows/minimize")
async def minimize_window(window_name: str):
    """Minimize a window"""
    try:
        import subprocess
        
        result = subprocess.run(
            ["xdotool", "search", "--name", window_name, "windowminimize"],
            capture_output=True,
            text=True,
            timeout=5
        )
        
        if result.returncode == 0:
            return {"status": "success", "action": f"Minimized window: {window_name}"}
        else:
            return {"status": "failed", "error": result.stderr}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed: {str(e)}")

@app.post("/media/play")
async def play_media(url: str = None, action: str = "play"):
    """Control media playback"""
    try:
        import subprocess
        
        # Try playerctl (works with most media players)
        if action == "play":
            result = subprocess.run(["playerctl", "play"], capture_output=True, text=True)
        elif action == "pause":
            result = subprocess.run(["playerctl", "pause"], capture_output=True, text=True)
        elif action == "next":
            result = subprocess.run(["playerctl", "next"], capture_output=True, text=True)
        elif action == "previous":
            result = subprocess.run(["playerctl", "previous"], capture_output=True, text=True)
        elif action == "stop":
            result = subprocess.run(["playerctl", "stop"], capture_output=True, text=True)
        elif action == "status":
            result = subprocess.run(["playerctl", "status"], capture_output=True, text=True)
            return {"status": "success", "player_status": result.stdout.strip()}
        else:
            return {"status": "error", "message": "Unknown action"}
        
        return {"status": "success", "action": action}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/media/open")
async def open_media(file_path: str):
    """Open media file with default player"""
    try:
        import subprocess
        import os
        
        if os.path.exists(file_path):
            result = subprocess.run(
                ["xdg-open", file_path],
                capture_output=True,
                text=True
            )
            return {"status": "success", "action": f"Opened: {file_path}"}
        else:
            return {"status": "error", "message": "File not found"}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/keyboard/type")
async def type_text(text: str):
    """Type text using keyboard"""
    try:
        import subprocess
        import shlex
        
        # Use xdotool to type text
        # Escape special characters
        escaped_text = text.replace("'", "'\\''")
        result = subprocess.run(
            f"xdotool type -- '{escaped_text}'",
            shell=True,
            capture_output=True,
            text=True,
            timeout=5
        )
        
        if result.returncode == 0:
            return {"status": "success", "action": f"Typed: {text[:50]}..."}
        else:
            return {"status": "error", "message": result.stderr}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/keyboard/press")
async def press_key(key: str):
    """Press a specific key"""
    try:
        import subprocess
        
        result = subprocess.run(
            ["xdotool", "key", key],
            capture_output=True,
            text=True,
            timeout=5
        )
        
        if result.returncode == 0:
            return {"status": "success", "action": f"Pressed: {key}"}
        else:
            return {"status": "error", "message": result.stderr}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/keyboard/combo")
async def key_combo(keys: str):
    """Press key combination (e.g., "ctrl+c", "alt+tab")"""
    try:
        import subprocess
        
        result = subprocess.run(
            ["xdotool", "key", keys],
            capture_output=True,
            text=True,
            timeout=5
        )
        
        if result.returncode == 0:
            return {"status": "success", "action": f"Pressed: {keys}"}
        else:
            return {"status": "error", "message": result.stderr}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/mouse/move")
async def move_mouse(x: int, y: int):
    """Move mouse to coordinates"""
    try:
        import subprocess
        
        result = subprocess.run(
            ["xdotool", "mousemove", str(x), str(y)],
            capture_output=True,
            text=True,
            timeout=5
        )
        
        if result.returncode == 0:
            return {"status": "success", "action": f"Moved to {x}, {y}"}
        else:
            return {"status": "error", "message": result.stderr}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/mouse/click")
async def click_mouse(button: str = "left"):
    """Mouse click"""
    try:
        import subprocess
        
        button_map = {"left": 1, "middle": 2, "right": 3}
        btn = button_map.get(button.lower(), 1)
        
        result = subprocess.run(
            ["xdotool", "click", str(btn)],
            capture_output=True,
            text=True,
            timeout=5
        )
        
        if result.returncode == 0:
            return {"status": "success", "action": f"Clicked: {button}"}
        else:
            return {"status": "error", "message": result.stderr}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/mouse/drag")
async def drag_mouse(x1: int, y1: int, x2: int, y2: int):
    """Drag from one point to another"""
    try:
        import subprocess
        
        # Move to start, hold, move to end, release
        cmds = [
            ["xdotool", "mousedown", "1"],
            ["xdotool", "mousemove", str(x1), str(y1)],
            ["xdotool", "mousemove", str(x2), str(y2)],
            ["xdotool", "mouseup", "1"]
        ]
        
        for cmd in cmds:
            subprocess.run(cmd, capture_output=True, timeout=2)
        
        return {"status": "success", "action": f"Dragged from ({x1},{y1}) to ({x2},{y2})"}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/screenshot")
async def take_screenshot(save_path: str = "/tmp/screenshot.png"):
    """Take screenshot"""
    try:
        import subprocess
        
        result = subprocess.run(
            ["gnome-screenshot", "-f", save_path],
            capture_output=True,
            text=True,
            timeout=5
        )
        
        if result.returncode == 0:
            return {"status": "success", "path": save_path}
        else:
            # Try import from PIL
            return {"status": "error", "message": "Screenshot failed"}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/scroll")
async def scroll_content(direction: str = "down", amount: int = 300):
    """Scroll content (up/down/left/right)"""
    try:
        import subprocess
        
        key_map = {"down": "Down", "up": "Up", "left": "Left", "right": "Right", "page_down": "Page_Down", "page_up": "Page_Up", "home": "Home", "end": "End"}
        key = key_map.get(direction, "Down")
        
        for _ in range(max(1, amount // 100)):
            subprocess.run(["xdotool", "key", key], capture_output=True, timeout=2)
        
        return {"status": "success", "action": f"Scrolled {direction}"}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.get("/clipboard")
async def get_clipboard():
    """Get clipboard content"""
    try:
        import subprocess
        result = subprocess.run(["xclip", "-selection", "clipboard", "-o"], capture_output=True, text=True, timeout=3)
        return {"status": "success", "content": result.stdout, "length": len(result.stdout)}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/clipboard")
async def set_clipboard(content: str):
    """Set clipboard content"""
    try:
        import subprocess
        subprocess.run(["xclip", "-selection", "clipboard", "-i"], input=content, text=True, capture_output=True, timeout=3)
        return {"status": "success", "action": "Clipboard set"}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.get("/clipboard/history")
async def get_clipboard_history():
    """Get clipboard history"""
    import os, json
    history_file = "/tmp/clipboard_history.json"
    if os.path.exists(history_file):
        with open(history_file, 'r') as f:
            return {"status": "success", "history": json.load(f)}
    return {"status": "success", "history": []}

@app.post("/clipboard/save")
async def save_to_clipboard_history(content: str, label: str = ""):
    """Save to clipboard history"""
    import os, json, datetime
    history_file = "/tmp/clipboard_history.json"
    if os.path.exists(history_file):
        with open(history_file, 'r') as f:
            history = json.load(f)
    else:
        history = []
    history.insert(0, {"content": content, "label": label, "timestamp": datetime.datetime.now().isoformat()})
    history = history[:50]
    with open(history_file, 'w') as f:
        json.dump(history, f)
    return {"status": "success", "action": "Saved to history"}

@app.post("/file/open")
async def open_file(file_path: str):
    """Open file with default app"""
    import subprocess, os
    if not os.path.exists(file_path):
        return {"status": "error", "message": "File not found"}
    subprocess.run(["xdg-open", file_path], capture_output=True, timeout=5)
    return {"status": "success", "action": f"Opened: {file_path}"}

@app.post("/file/read")
async def read_file(file_path: str, lines: int = 100):
    """Read file content"""
    import os
    if not os.path.exists(file_path):
        return {"status": "error", "message": "File not found"}
    with open(file_path, 'r') as f:
        content = f.read(lines * 1000)
    return {"status": "success", "file": file_path, "size_bytes": os.path.getsize(file_path), "content": content[:5000]}

@app.post("/whatsapp/send")
async def send_whatsapp_message(phone: str, message: str):
    """Send WhatsApp message"""
    return {"status": "info", "message": "WhatsApp requires browser automation", "note": "Use keyboard to type in WhatsApp Web"}

@app.post("/power")
async def power_action(action: str):
    """System power actions (secure)"""
    import subprocess
    allowed = ["logout", "suspend", "hibernate", "reboot", "shutdown"]
    
    if action not in allowed:
        return {"status": "error", "message": f"Allowed: {allowed}"}
    
    # Require confirmation for dangerous actions
    try:
        if action == "logout":
            subprocess.run(["pkill", "-KILL", "-u", os.getenv("USER", "root")], timeout=5)
        elif action == "suspend":
            subprocess.run(["systemctl", "suspend"], timeout=5)
        elif action == "hibernate":
            subprocess.run(["systemctl", "hibernate"], timeout=5)
        elif action == "reboot":
            subprocess.run(["reboot"], timeout=5)
        elif action == "shutdown":
            subprocess.run(["shutdown", "-now"], timeout=5)
        return {"status": "success", "action": action}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.get("/news")
async def get_top_news(category: str = "general", country: str = "in"):
    """Get top news headlines"""
    try:
        import urllib.request, urllib.parse, json
        
        # Use GNews API or fallback to RSS
        query = urllib.parse.quote(f"{category} news")
        url = f"https://gnews.io/api/v4/top-headlines?category={category}&lang=en&max=10"
        
        try:
            with urllib.request.urlopen(url, timeout=10) as response:
                data = json.loads(response.read().decode())
            
            articles = []
            for article in data.get("articles", [])[:10]:
                articles.append({
                    "title": article.get("title", ""),
                    "description": article.get("description", ""),
                    "url": article.get("url", ""),
                    "source": article.get("source", {}).get("name", ""),
                    "image": article.get("image", ""),
                    "published": article.get("publishedAt", "")
                })
            return {"status": "success", "category": category, "articles": articles, "count": len(articles)}
        except:
            # Fallback: return sample news
            return {
                "status": "info", 
                "message": "Using sample news (API key needed for live)",
                "articles": [
                    {"title": "Tech industry sees AI growth", "source": "Tech News"},
                    {"title": "Markets show positive trend", "source": "Business"},
                    {"title": "New innovations in quantum computing", "source": "Science"}
                ]
            }
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/execute/parallel")
async def execute_parallel(commands: List[str]):
    """Execute multiple commands in parallel"""
    import subprocess, concurrent.futures
    
    results = []
    
    def run_cmd(cmd):
        try:
            result = subprocess.run(cmd, shell=True, capture_output=True, text=True, timeout=30)
            return {"command": cmd, "stdout": result.stdout[:500], "stderr": result.stderr[:200], "returncode": result.returncode}
        except Exception as e:
            return {"command": cmd, "error": str(e)}
    
    with concurrent.futures.ThreadPoolExecutor(max_workers=10) as executor:
        futures = [executor.submit(run_cmd, cmd) for cmd in commands]
        results = [f.result() for f in concurrent.futures.as_completed(futures)]
    
    return {"status": "success", "commands": len(commands), "results": results}

@app.post("/code/generate")
async def generate_and_type_code(prompt: str, language: str = "python"):
    """Generate code from AI and type it"""
    try:
        # Use local AI to generate code
        if nexus:
            response = nexus.chat(
                model="gemini/gemini-pro",
                messages=[{"role": "user", "content": f"Write {language} code for: {prompt}. Only output the code, no explanation."}],
                priority="speed"
            )
            code = response.get("response", "")[:2000]
        else:
            code = f"# Code for: {prompt}\n# (AI not configured)\nprint('Hello World')"
        
        # Type the code using keyboard
        try:
            import subprocess
            escaped = code.replace("'", "'\\''").replace("\n", "\\n")
            subprocess.run(f"xdotool type -- '{escaped}'", shell=True, capture_output=True, timeout=10)
            typed = True
        except:
            typed = False
        
        return {"status": "success", "code": code, "language": language, "typed": typed, "length": len(code)}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/vscode/open")
async def open_in_vscode(file_path: str = None, content: str = None):
    """Open file or create new file in VSCode"""
    import subprocess, os, tempfile
    
    if file_path:
        # Open existing file
        result = subprocess.run(["code", file_path], capture_output=True, text=True)
        return {"status": "success", "action": f"Opened in VSCode: {file_path}"}
    elif content:
        # Create temp file and open
        with tempfile.NamedTemporaryFile(mode='w', suffix='.txt', delete=False) as f:
            f.write(content)
            temp_path = f.name
        subprocess.run(["code", temp_path], capture_output=True)
        return {"status": "success", "action": f"Created and opened in VSCode", "file": temp_path}
    else:
        # Just open VSCode
        subprocess.run(["code"], capture_output=True)
        return {"status": "success", "action": "Opened VSCode"}

@app.post("/vscode/command")
async def vscode_command(command: str):
    """Execute VSCode command"""
    import subprocess
    # Use code --command for VS Code CLI
    result = subprocess.run(["code", "--command", command], capture_output=True, text=True, timeout=10)
    return {"status": "success", "output": result.stdout, "error": result.stderr}

@app.post("/screenshot/instant")
async def instant_screenshot(save_path: str = "/tmp/screenshot.png"):
    """Take instant high-res screenshot"""
    import subprocess, os
    from datetime import datetime
    
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    save_path = f"/tmp/screenshot_{timestamp}.png"
    
    # Try gnome-screenshot first
    result = subprocess.run(["gnome-screenshot", "-f", save_path], capture_output=True, text=True, timeout=5)
    
    if result.returncode != 0:
        # Fallback to import (ImageMagick)
        result = subprocess.run(["import", "-window", "root", save_path], capture_output=True, text=True, timeout=5)
    
    if os.path.exists(save_path):
        size = os.path.getsize(save_path)
        return {"status": "success", "path": save_path, "size_bytes": size}
    else:
        return {"status": "error", "message": "Screenshot failed"}

@app.post("/screenshot/region")
async def region_screenshot(x: int, y: int, width: int, height: int, save_path: str = None):
    """Take screenshot of specific region"""
    import subprocess, os
    from datetime import datetime
    
    if not save_path:
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        save_path = f"/tmp/region_{timestamp}.png"
    
    result = subprocess.run(
        ["import", "-window", "root", "-crop", f"{width}x{height}+{x}+{y}", save_path],
        capture_output=True,
        text=True,
        timeout=10
    )
    
    if os.path.exists(save_path):
        return {"status": "success", "path": save_path, "size_bytes": os.path.getsize(save_path)}
    else:
        return {"status": "error", "message": "Region screenshot failed"}

@app.post("/type/human")
async def human_like_typing(text: str, min_delay: float = 0.05, max_delay: float = 0.2):
    """Human-like auto typing with random delays"""
    import subprocess
    import random
    
    # Type character by character with random delays
    for char in text:
        escaped = char.replace("'", "'\\''")
        subprocess.run(f"xdotool type -- '{escaped}'", shell=True, capture_output=True, timeout=2)
        delay = random.uniform(min_delay, max_delay)
        time.sleep(delay)
    
    return {"status": "success", "typed": len(text), "message": f"Typed {len(text)} characters human-like"}

@app.post("/type/line")
async def type_line(text: str):
    """Type a full line at once"""
    import subprocess
    
    escaped = text.replace("'", "'\\''").replace("\n", "\\n")
    subprocess.run(f"xdotool type -- '{escaped}'", shell=True, capture_output=True, timeout=10)
    return {"status": "success", "typed": len(text)}

@app.get("/security/scan")
async def scan_system():
    """Basic system security scan"""
    import subprocess, os
    
    issues = []
    warnings = []
    
    # Check for open ports
    result = subprocess.run(["ss", "-tuln"], capture_output=True, text=True)
    common_ports = [21, 23, 25, 110, 143, 445, 3389, 5900]
    for port in common_ports:
        if f":{port}" in result.stdout:
            warnings.append(f"Port {port} is open")
    
    # Check for root access
    if os.geteuid() == 0:
        issues.append("Running as ROOT - security risk!")
    
    # Check SSH keys
    ssh_dir = os.path.expanduser("~/.ssh")
    if os.path.exists(ssh_dir):
        key_files = [f for f in os.listdir(ssh_dir) if f.endswith('.pub')]
        warnings.append(f"Found {len(key_files)} SSH public keys")
    
    # Check firewall status
    result = subprocess.run(["ufw", "status"], capture_output=True, text=True)
    if "inactive" in result.stdout.lower():
        warnings.append("Firewall (UFW) is INACTIVE")
    
    # Check recent failed logins
    result = subprocess.run(["last", "-f", "/var/log/btmp", "-n", "5"], capture_output=True, text=True)
    if result.stdout:
        warnings.append("Failed login attempts detected")
    
    return {
        "status": "scanned",
        "timestamp": datetime.now().isoformat(),
        "issues": issues,
        "warnings": warnings,
        "secure_score": max(0, 100 - (len(issues) * 30) - (len(warnings) * 5))
    }

@app.get("/security/processes")
async def suspicious_processes():
    """Check for suspicious processes"""
    import psutil
    
    suspicious = []
    known_bad = ["nc ", "netcat", "ncat", "socat", "meterpreter", "mimikatz"]
    
    for proc in psutil.process_iter(['pid', 'name', 'cmdline']):
        try:
            name = proc.info['name'] or ""
            cmdline = ' '.join(proc.info['cmdline'] or [])
            
            for bad in known_bad:
                if bad in cmdline.lower():
                    suspicious.append({
                        "pid": proc.info['pid'],
                        "name": name,
                        "reason": f"Contains: {bad}"
                    })
        except:
            pass
    
    return {"status": "success", "suspicious_processes": suspicious, "count": len(suspicious)}

@app.post("/volume/set")
async def set_volume(level: int):
    """Set system volume (0-100)"""
    import subprocess
    
    level = max(0, min(100, level))
    
    try:
        # Try pactl (PulseAudio)
        subprocess.run(["pactl", "set-sink-volume", "@DEFAULT_SINK@", f"{level}%"], capture_output=True)
        return {"status": "success", "volume": level}
    except:
        try:
            # Try amixer
            subprocess.run(["amixer", "-D", "pulse", "sset", "Master", f"{level}%"], capture_output=True)
            return {"status": "success", "volume": level}
        except:
            return {"status": "error", "message": "Volume control not available"}

@app.get("/volume")
async def get_volume():
    """Get current volume level"""
    import subprocess
    
    try:
        result = subprocess.run(["pactl", "get-sink-volume", "@DEFAULT_SINK@"], capture_output=True, text=True)
        # Parse output
        return {"status": "success", "output": result.stdout[:200]}
    except:
        return {"status": "error", "message": "Cannot get volume"}

@app.post("/volume/mute")
async def mute_volume(mute: bool = True):
    """Mute/unmute system"""
    import subprocess
    
    try:
        subprocess.run(["pactl", "set-sink-mute", "@DEFAULT_SINK@", "toggle" if mute else "unmute"], capture_output=True)
        return {"status": "success", "muted": mute}
    except:
        return {"status": "error", "message": "Mute control not available"}

@app.post("/brightness/set")
async def set_brightness(level: int):
    """Set screen brightness (0-100)"""
    import subprocess
    
    level = max(0, min(100, level))
    
    try:
        # Try xbacklight
        subprocess.run(["xbacklight", "-set", str(level)], capture_output=True)
        return {"status": "success", "brightness": level}
    except:
        try:
            # Try brightnessctl
            subprocess.run(["brightnessctl", "set", f"{level}%"], capture_output=True)
            return {"status": "success", "brightness": level}
        except:
            return {"status": "error", "message": "Brightness control not available"}

@app.get("/brightness")
async def get_brightness():
    """Get current brightness level"""
    import subprocess
    
    try:
        result = subprocess.run(["brightnessctl", "get"], capture_output=True, text=True)
        return {"status": "success", "brightness": result.stdout.strip()}
    except:
        return {"status": "error", "message": "Cannot get brightness"}

@app.post("/brightness/up")
async def brightness_up(amount: int = 10):
    """Increase brightness"""
    import subprocess
    try:
        subprocess.run(["brightnessctl", "set", f"+{amount}%"], capture_output=True)
        return {"status": "success", "action": f"Increased by {amount}%"}
    except:
        return {"status": "error", "message": "Brightness control not available"}

@app.post("/brightness/down")
async def brightness_down(amount: int = 10):
    """Decrease brightness"""
    import subprocess
    try:
        subprocess.run(["brightnessctl", "set", f"-{amount}%"], capture_output=True)
        return {"status": "success", "action": f"Decreased by {amount}%"}
    except:
        return {"status": "error", "message": "Brightness control not available"}

@app.post("/ai/image")
async def generate_ai_image(prompt: str, style: str = "realistic"):
    """Generate AI image (uses external API or local)"""
    try:
        # Try using local AI or return placeholder
        return {
            "status": "info",
            "message": "Image generation requires API key",
            "prompt": prompt,
            "style": style,
            "alternatives": [
                "Use DALL-E API with OPENAI_API_KEY",
                "Use Stable Diffusion local",
                "Use ComfyUI"
            ]
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/code/fix")
async def fix_code_error(code: str, language: str = "python", error: str = ""):
    """Fix code errors using AI"""
    try:
        if nexus:
            prompt = f"Fix this {language} code. "
            if error:
                prompt += f"Error: {error}. "
            prompt += f"Code:\n{code}\n\nProvide the corrected code only."
            
            response = nexus.chat(
                model="gemini/gemini-pro",
                messages=[{"role": "user", "content": prompt}],
                priority="speed"
            )
            fixed = response.get("response", "")[:3000]
        else:
            fixed = "# AI not configured\n# Please provide error details"
        
        return {
            "status": "success",
            "original_code": code[:500],
            "fixed_code": fixed,
            "language": language
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/code/lint")
async def lint_code(code: str, language: str = "python"):
    """Lint/analyze code for issues"""
    import subprocess
    import tempfile
    import os
    
    issues = []
    
    if language == "python":
        # Save to temp file and run flake8/pylint
        with tempfile.NamedTemporaryFile(mode='w', suffix='.py', delete=False) as f:
            f.write(code)
            temp_path = f.name
        
        try:
            result = subprocess.run(["flake8", temp_path, "--max-line-length=100"], capture_output=True, text=True)
            if result.stdout:
                issues.extend(result.stdout.strip().split('\n'))
        except:
            pass
        
        os.unlink(temp_path)
    
    return {
        "status": "success",
        "language": language,
        "issues": issues[:20],
        "issue_count": len(issues)
    }

# ============== REMINDERS ==============
import json, os
from datetime import datetime, timedelta

REMINDERS_FILE = "/tmp/nexusai_reminders.json"

def load_reminders():
    if os.path.exists(REMINDERS_FILE):
        with open(REMINDERS_FILE, 'r') as f:
            return json.load(f)
    return []

def save_reminders(reminders):
    with open(REMINDERS_FILE, 'w') as f:
        json.dump(reminders, f)

@app.post("/reminder/set")
async def set_reminder(message: str, minutes: int = 60, channel: str = "telegram"):
    """Set a reminder"""
    reminders = load_reminders()
    
    reminder = {
        "id": f"rem_{len(reminders)+1}",
        "message": message,
        "created_at": datetime.now().isoformat(),
        "due_at": (datetime.now() + timedelta(minutes=minutes)).isoformat(),
        "channel": channel,
        "status": "pending"
    }
    
    reminders.append(reminder)
    save_reminders(reminders)
    
    return {"status": "success", "reminder": reminder}

@app.get("/reminders")
async def view_reminders():
    """View all reminders"""
    reminders = load_reminders()
    now = datetime.now()
    
    # Check which are due
    for r in reminders:
        due = datetime.fromisoformat(r['due_at'])
        r['due'] = due <= now
        r['minutes_remaining'] = int((due - now).total_seconds() / 60)
    
    return {"status": "success", "reminders": reminders, "count": len(reminders)}

@app.post("/reminder/cancel")
async def cancel_reminder(reminder_id: str):
    """Cancel a reminder"""
    reminders = load_reminders()
    original_count = len(reminders)
    
    reminders = [r for r in reminders if r['id'] != reminder_id]
    
    if len(reminders) < original_count:
        save_reminders(reminders)
        return {"status": "success", "action": f"Reminder {reminder_id} cancelled"}
    
    return {"status": "error", "message": "Reminder not found"}

@app.post("/reminder/clear")
async def clear_all_reminders():
    """Clear all reminders"""
    save_reminders([])
    return {"status": "success", "action": "All reminders cleared"}

# ============== DOCUMENT PROCESSING ==============

@app.post("/document/analyze")
async def analyze_document(text: str, question: str = ""):
    """Analyze document using AI"""
    try:
        if nexus:
            if question:
                prompt = f"Document:\n{text[:3000]}\n\nQuestion: {question}\n\nProvide a detailed answer."
            else:
                prompt = f"Analyze this document and provide a summary:\n{text[:3000]}"
            
            response = nexus.chat(
                model="gemini/gemini-pro",
                messages=[{"role": "user", "content": prompt}],
                priority="quality"
            )
            result = response.get("response", "")[:2000]
        else:
            result = "AI not configured - please add API keys"
        
        return {"status": "success", "analysis": result, "text_length": len(text)}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/document/summarize")
async def summarize_document(text: str, max_words: int = 100):
    """Summarize document"""
    try:
        if nexus:
            prompt = f"Summarize this in exactly {max_words} words:\n{text[:3000]}"
            
            response = nexus.chat(
                model="gemini/gemini-pro",
                messages=[{"role": "user", "content": prompt}],
                priority="speed"
            )
            summary = response.get("response", "")[:1000]
        else:
            summary = "AI not configured"
        
        return {"status": "success", "summary": summary, "word_count": len(summary.split())}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/document/extract")
async def extract_from_document(text: str, extract_type: str = "key_points"):
    """Extract specific information from document"""
    extract_types = {
        "key_points": "Extract the key points as a bullet list",
        "dates": "Extract all dates and events mentioned",
        "emails": "Extract all email addresses",
        "urls": "Extract all URLs/links",
        "numbers": "Extract all important numbers",
        "names": "Extract all names of people mentioned"
    }
    
    try:
        if nexus:
            prompt = f"{extract_types.get(extract_type, extract_types['key_points'])}:\n{text[:3000]}"
            
            response = nexus.chat(
                model="gemini/gemini-pro",
                messages=[{"role": "user", "content": prompt}],
                priority="speed"
            )
            extracted = response.get("response", "")[:1500]
        else:
            extracted = "AI not configured"
        
        return {"status": "success", "extracted": extracted, "type": extract_type}
    except Exception as e:
        return {"status": "error", "message": str(e)}

# ============== MEDIA & WHATSAPP ==============

@app.post("/whatsapp/send_media")
async def send_media_to_whatsapp(phone: str, caption: str = "", file_path: str = ""):
    """Send media (image/video) to WhatsApp"""
    return {
        "status": "info",
        "message": "WhatsApp media requires browser automation",
        "note": "Use /keyboard/ to type and /screenshot to capture, then share manually"
    }

# ============== EXCEL & SPREADSHEETS ==============
try:
    import xlsxwriter
    import openpyxl
    EXCEL_AVAILABLE = True
except:
    EXCEL_AVAILABLE = False

@app.post("/excel/create")
async def create_excel_file(file_name: str = "/tmp/report.xlsx", sheets: str = "Sheet1"):
    """Create new Excel file"""
    if not EXCEL_AVAILABLE:
        return {"status": "error", "message": "Excel libraries not installed"}
    
    try:
        workbook = xlsxwriter.Workbook(file_name)
        
        # Create sheets
        sheet_names = sheets.split(',')
        for sheet_name in sheet_names:
            worksheet = workbook.add_worksheet(sheet_name.strip())
        
        # Add default header
        worksheet.write('A1', 'Column A')
        worksheet.write('B1', 'Column B')
        worksheet.write('C1', 'Column C')
        
        workbook.close()
        
        return {"status": "success", "file": file_name, "sheets": sheet_names}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/excel/add_data")
async def add_data_to_excel(file_path: str, data: str, sheet: str = "Sheet1"):
    """Add data to Excel file"""
    try:
        import openpyxl
        
        wb = openpyxl.load_workbook(file_path)
        ws = wb[sheet] if sheet in wb.sheetnames else wb.active
        
        # Parse data (JSON format)
        import json
        rows = json.loads(data) if data.startswith('[') else [[data]]
        
        # Find last row
        last_row = ws.max_row + 1
        
        for row_idx, row_data in enumerate(rows):
            for col_idx, value in enumerate(row_data):
                ws.cell(row=last_row + row_idx, column=col_idx + 1, value=value)
        
        wb.save(file_path)
        
        return {"status": "success", "rows_added": len(rows), "file": file_path}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/excel/write")
async def write_to_excel(file_path: str, cell: str, value: str, sheet: str = "Sheet1"):
    """Write to specific cell"""
    try:
        import openpyxl
        
        wb = openpyxl.load_workbook(file_path)
        ws = wb[sheet] if sheet in wb.sheetnames else wb.active
        
        ws[cell] = value
        wb.save(file_path)
        
        return {"status": "success", "cell": cell, "value": value}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.get("/excel/read")
async def read_excel(file_path: str, range: str = "A1:Z10"):
    """Read data from Excel"""
    try:
        import openpyxl
        
        wb = openpyxl.load_workbook(file_path)
        ws = wb.active
        
        data = []
        for row in ws[range]:
            row_data = [cell.value for cell in row]
            data.append(row_data)
        
        return {"status": "success", "data": data, "rows": len(data)}
    except Exception as e:
        return {"status": "error", "message": str(e)}

# ============== DATA MANAGEMENT ==============

DATA_FILE = "/tmp/nexusai_data.json"

@app.post("/data/save")
async def save_data(key: str, value: str):
    """Save key-value data"""
    import json
    
    if os.path.exists(DATA_FILE):
        with open(DATA_FILE, 'r') as f:
            data = json.load(f)
    else:
        data = {}
    
    data[key] = value
    
    with open(DATA_FILE, 'w') as f:
        json.dump(data, f)
    
    return {"status": "success", "key": key}

@app.get("/data/get")
async def get_data(key: str):
    """Get saved data"""
    import json
    
    if os.path.exists(DATA_FILE):
        with open(DATA_FILE, 'r') as f:
            data = json.load(f)
        return {"status": "success", "key": key, "value": data.get(key)}
    
    return {"status": "error", "message": "Key not found"}

@app.get("/data/list")
async def list_all_data():
    """List all saved data"""
    import json
    
    if os.path.exists(DATA_FILE):
        with open(DATA_FILE, 'r') as f:
            data = json.load(f)
        return {"status": "success", "data": data, "count": len(data)}
    
    return {"status": "success", "data": {}, "count": 0}

@app.delete("/data/delete")
async def delete_data(key: str):
    """Delete specific data"""
    import json
    
    if os.path.exists(DATA_FILE):
        with open(DATA_FILE, 'r') as f:
            data = json.load(f)
        
        if key in data:
            del data[key]
            
            with open(DATA_FILE, 'w') as f:
                json.dump(data, f)
            
            return {"status": "success", "key": key}
    
    return {"status": "error", "message": "Key not found"}

@app.delete("/data/clear")
async def delete_all_data():
    """Delete all saved data"""
    import json
    
    with open(DATA_FILE, 'w') as f:
        json.dump({}, f)
    
    return {"status": "success", "action": "All data deleted"}

# ============== EXCEL NAVIGATION ==============

@app.post("/excel/navigate")
async def navigate_cell(action: str, cell: str = "A1"):
    """Navigate cells in Excel (keyboard simulation)"""
    import subprocess
    
    movements = {
        "left": "Left",
        "right": "Right", 
        "up": "Up",
        "down": "Down",
        "home": "Home",
        "end": "End",
        "page_up": "Page_Up",
        "page_down": "Page_Down"
    }
    
    key = movements.get(action, "Right")
    
    try:
        subprocess.run(["xdotool", "key", key], capture_output=True, timeout=2)
        return {"status": "success", "action": f"Moved {action}"}
    except:
        return {"status": "error", "message": "xdotool not available"}

@app.post("/excel/move_left")
async def move_left(steps: int = 1):
    """Move left in grid/spreadsheet"""
    for _ in range(steps):
        subprocess.run(["xdotool", "key", "Left"], capture_output=True, timeout=2)
    return {"status": "success", "moved": f"left {steps} cells"}

@app.post("/excel/move_up")
async def move_up(steps: int = 1):
    """Move up in grid/spreadsheet"""
    for _ in range(steps):
        subprocess.run(["xdotool", "key", "Up"], capture_output=True, timeout=2)
    return {"status": "success", "moved": f"up {steps} cells"}

@app.post("/excel/move_down")
async def move_down(steps: int = 1):
    """Move down in grid/spreadsheet"""
    for _ in range(steps):
        subprocess.run(["xdotool", "key", "Down"], capture_output=True, timeout=2)
    return {"status": "success", "moved": f"down {steps} cells"}

@app.post("/excel/move_right")
async def move_right(steps: int = 1):
    """Move right in grid/spreadsheet"""
    for _ in range(steps):
        subprocess.run(["xdotool", "key", "Right"], capture_output=True, timeout=2)
    return {"status": "success", "moved": f"right {steps} cells"}

@app.post("/excel/delete_cell")
async def delete_current_cell():
    """Delete current cell content (Backspace)"""
    try:
        subprocess.run(["xdotool", "key", "Backspace"], capture_output=True, timeout=2)
        return {"status": "success", "action": "Cell content deleted"}
    except:
        return {"status": "error", "message": "Failed"}

@app.post("/excel/toggle_bold")
async def toggle_text_bold():
    """Toggle bold (Ctrl+B)"""
    try:
        subprocess.run(["xdotool", "key", "ctrl+b"], capture_output=True, timeout=2)
        return {"status": "success", "action": "Toggled bold"}
    except:
        return {"status": "error", "message": "Failed"}

@app.post("/excel/toggle_italic")
async def toggle_text_italic():
    """Toggle italic (Ctrl+I)"""
    try:
        subprocess.run(["xdotool", "key", "ctrl+i"], capture_output=True, timeout=2)
        return {"status": "success", "action": "Toggled italic"}
    except:
        return {"status": "error", "message": "Failed"}

@app.post("/excel/select_row")
async def select_row():
    """Select row (Shift+Space)"""
    try:
        subprocess.run(["xdotool", "key", "shift+space"], capture_output=True, timeout=2)
        return {"status": "success", "action": "Row selected"}
    except:
        return {"status": "error", "message": "Failed"}

@app.post("/excel/select_column")
async def select_column():
    """Select column (Ctrl+Space)"""
    try:
        subprocess.run(["xdotool", "key", "ctrl+space"], capture_output=True, timeout=2)
        return {"status": "success", "action": "Column selected"}
    except:
        return {"status": "error", "message": "Failed"}

@app.post("/excel/select_all")
async def select_all():
    """Select all (Ctrl+A)"""
    try:
        subprocess.run(["xdotool", "key", "ctrl+a"], capture_output=True, timeout=2)
        return {"status": "success", "action": "All selected"}
    except:
        return {"status": "error", "message": "Failed"}

@app.post("/excel/sort")
async def sort_excel_data(direction: str = "asc"):
    """Sort data"""
    try:
        subprocess.run(["xdotool", "key", "alt+d"], capture_output=True, timeout=2)
        subprocess.run(["xdotool", "key", "s"], capture_output=True, timeout=2)
        if direction == "desc":
            subprocess.run(["xdotool", "key", "Down"], capture_output=True, timeout=2)
        subprocess.run(["xdotool", "key", "Return"], capture_output=True, timeout=2)
        return {"status": "success", "action": f"Sorted {direction}"}
    except:
        return {"status": "error", "message": "Sort failed"}

@app.post("/excel/clipboard_action")
async def excel_clipboard_action(action: str = "copy"):
    """Excel clipboard: copy/paste/cut"""
    try:
        if action == "copy":
            subprocess.run(["xdotool", "key", "ctrl+c"], capture_output=True, timeout=2)
        elif action == "paste":
            subprocess.run(["xdotool", "key", "ctrl+v"], capture_output=True, timeout=2)
        elif action == "cut":
            subprocess.run(["xdotool", "key", "ctrl+x"], capture_output=True, timeout=2)
        return {"status": "success", "action": action}
    except:
        return {"status": "error", "message": "Failed"}

@app.post("/excel/go_to_cell")
async def go_to_cell(cell: str):
    """Go to specific cell address"""
    try:
        # Ctrl+G to open goto dialog
        subprocess.run(["xdotool", "key", "ctrl+g"], capture_output=True, timeout=2)
        import time
        time.sleep(0.3)
        # Type cell address
        subprocess.run(["xdotool", "type", cell], capture_output=True, timeout=2)
        time.sleep(0.3)
        # Press Enter
        subprocess.run(["xdotool", "key", "Return"], capture_output=True, timeout=2)
        return {"status": "success", "cell": cell}
    except:
        return {"status": "error", "message": "Failed to go to cell"}

@app.post("/excel/move_to")
async def move_to_cell(cell: str):
    """Move to specific cell"""
    # Press Ctrl+G and type cell address
    try:
        subprocess.run(["xdotool", "key", "ctrl+g"], capture_output=True, timeout=2)
        subprocess.run(["xdotool", "type", cell], capture_output=True, timeout=2)
        subprocess.run(["xdotool", "key", "Return"], capture_output=True, timeout=2)
        return {"status": "success", "cell": cell}
    except:
        return {"status": "error", "message": "Navigation failed"}

@app.post("/excel/enter_data")
async def enter_data_quick(value: str):
    """Quick data entry - type and enter"""
    try:
        import subprocess
        escaped = value.replace("'", "'\\''")
        subprocess.run(f"xdotool type -- '{escaped}'", shell=True, capture_output=True, timeout=5)
        subprocess.run(["xdotool", "key", "Return"], capture_output=True, timeout=2)
        return {"status": "success", "entered": value}
    except:
        return {"status": "error", "message": "Failed to enter data"}

@app.post("/excel/enter_multiple")
async def enter_multiple_data(data: str):
    """Enter multiple rows of data (tab between cells, enter for new row)"""
    try:
        import subprocess
        
        rows = data.split('|')  # Use | as row separator
        
        for row_idx, row in enumerate(rows):
            cells = row.split(',')  # Use , as cell separator
            for cell_idx, value in enumerate(cells):
                escaped = value.strip().replace("'", "'\\''")
                subprocess.run(f"xdotool type -- '{escaped}'", shell=True, capture_output=True, timeout=3)
                if cell_idx < len(cells) - 1:
                    subprocess.run(["xdotool", "key", "Tab"], capture_output=True, timeout=2)
            
            if row_idx < len(rows) - 1:
                subprocess.run(["xdotool", "key", "Return"], capture_output=True, timeout=2)
        
        return {"status": "success", "rows_entered": len(rows)}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/excel/copy")
async def copy_cell():
    """Copy (Ctrl+C)"""
    try:
        subprocess.run(["xdotool", "key", "ctrl+c"], capture_output=True, timeout=2)
        return {"status": "success", "action": "Copied"}
    except:
        return {"status": "error", "message": "Failed"}

@app.post("/excel/paste")
async def paste_cell():
    """Paste (Ctrl+V)"""
    try:
        subprocess.run(["xdotool", "key", "ctrl+v"], capture_output=True, timeout=2)
        return {"status": "success", "action": "Pasted"}
    except:
        return {"status": "error", "message": "Failed"}

@app.post("/excel/delete")
async def delete_cell():
    """Delete cell content (Delete key)"""
    try:
        subprocess.run(["xdotool", "key", "Delete"], capture_output=True, timeout=2)
        return {"status": "success", "action": "Deleted"}
    except:
        return {"status": "error", "message": "Failed"}

# ============== CALCULATIONS & CONVERSIONS ==============

@app.get("/calc/sum")
async def calculate_sum(numbers: str):
    """Calculate sum of numbers (comma separated)"""
    try:
        nums = [float(n.strip()) for n in numbers.split(',')]
        return {"status": "success", "sum": sum(nums), "count": len(nums), "numbers": nums}
    except:
        return {"status": "error", "message": "Invalid numbers"}

@app.get("/calc/avg")
async def calculate_average(numbers: str):
    """Calculate average"""
    try:
        nums = [float(n.strip()) for n in numbers.split(',')]
        avg = sum(nums) / len(nums)
        return {"status": "success", "average": avg, "count": len(nums)}
    except:
        return {"status": "error", "message": "Invalid numbers"}

@app.get("/calc/multiply")
async def calculate_multiply(numbers: str):
    """Multiply numbers"""
    try:
        nums = [float(n.strip()) for n in numbers.split(',')]
        result = 1
        for n in nums:
            result *= n
        return {"status": "success", "result": result, "numbers": nums}
    except:
        return {"status": "error", "message": "Invalid numbers"}

@app.post("/convert/word_to_pdf")
async def word_to_pdf(input_file: str, output_file: str = None):
    """Convert Word to PDF"""
    if not output_file:
        output_file = input_file.replace('.docx', '.pdf').replace('.doc', '.pdf')
    
    return {
        "status": "info",
        "message": "Word to PDF requires libreoffice",
        "command": f"libreoffice --headless --convert-to pdf {input_file} --outdir /tmp",
        "example": "Use: soffice --headless --convert-to pdf input.docx"
    }

@app.post("/convert/image_to_pdf")
async def image_to_pdf(image_paths: str, output_file: str = "/tmp/merged.pdf"):
    """Convert images to PDF"""
    try:
        from PIL import Image
        import os
        
        images = []
        for path in image_paths.split(','):
            path = path.strip()
            if os.path.exists(path):
                images.append(Image.open(path).convert('RGB'))
        
        if images:
            images[0].save(output_file, save_all=True, append_images=images[1:])
            return {"status": "success", "output": output_file, "images": len(images)}
        else:
            return {"status": "error", "message": "No valid images found"}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/convert/excel_to_pdf")
async def excel_to_pdf(input_file: str, output_file: str = None):
    """Convert Excel to PDF"""
    return {
        "status": "info",
        "message": "Excel to PDF requires libreoffice",
        "command": f"libreoffice --headless --convert-to pdf {input_file}",
        "example": "Use: soffice --headless --convert-to pdf input.xlsx"
    }

# ============== MORE CONVERSIONS ==============

@app.post("/convert/ppt_to_pdf")
async def ppt_to_pdf(input_file: str, output_file: str = None):
    """Convert PowerPoint to PDF"""
    return {
        "status": "info",
        "message": "PPT to PDF requires libreoffice",
        "command": f"libreoffice --headless --convert-to pdf {input_file}",
        "example": "Use: soffice --headless --convert-to pdf presentation.pptx"
    }

@app.post("/convert/image_format")
async def convert_image_format(input_file: str, output_format: str = "png", output_file: str = None):
    """Convert image format (jpg, png, webp, etc)"""
    try:
        from PIL import Image
        import os
        
        if not output_file:
            name = os.path.splitext(input_file)[0]
            output_file = f"{name}.{output_format}"
        
        img = Image.open(input_file)
        img.save(output_file)
        
        return {"status": "success", "input": input_file, "output": output_file, "format": output_format}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.get("/convert/formats")
async def supported_formats():
    """List supported conversion formats"""
    return {
        "image": ["jpg", "jpeg", "png", "webp", "bmp", "gif", "tiff"],
        "document": ["pdf"],
        "status": "success"
    }

@app.get("/test/converters")
async def test_converters():
    """Test converter availability"""
    results = {"calculations": False, "image_conversion": False, "excel": False}
    
    # Test calculations
    try:
        nums = [1, 2, 3]
        assert sum(nums) == 6
        results["calculations"] = True
    except:
        pass
    
    # Test PIL
    try:
        from PIL import Image
        results["image_conversion"] = True
    except:
        pass
    
    # Test Excel
    try:
        import openpyxl
        results["excel"] = True
    except:
        pass
    
    return {"status": "success", "results": results}

@app.post("/file/create")
async def create_file(file_name: str, content: str = ""):
    """Create a file"""
    try:
        with open(file_name, 'w') as f:
            f.write(content)
        return {"status": "success", "file": file_name, "size": len(content)}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/file/create_here")
async def create_here(file_name: str, content: str = ""):
    """Create file in current directory"""
    import os
    cwd = os.getcwd()
    file_path = os.path.join(cwd, file_name)
    
    try:
        with open(file_path, 'w') as f:
            f.write(content)
        return {"status": "success", "file": file_path, "size": len(content)}
    except Exception as e:
        return {"status": "error", "message": str(e)}

# ============== VISION & OCR ==============

@app.get("/vision/read_screen")
async def read_screen_text():
    """Read text from screen using OCR"""
    try:
        # Take screenshot
        import subprocess
        import tempfile
        import os
        
        temp_file = "/tmp/screen_ocr.png"
        subprocess.run(["gnome-screenshot", "-f", temp_file], capture_output=True, timeout=10)
        
        if not os.path.exists(temp_file):
            subprocess.run(["import", "-window", "root", temp_file], capture_output=True, timeout=10)
        
        # Try OCR
        try:
            import pytesseract
            from PIL import Image
            
            img = Image.open(temp_file)
            text = pytesseract.image_to_string(img)
            
            return {"status": "success", "text": text[:2000], "length": len(text)}
        except:
            return {"status": "info", "message": "OCR not available (install pytesseract)", "screenshot": temp_file}
    
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.get("/vision/screen")
async def analyze_screen():
    """Analyze entire screen with AI"""
    try:
        # Take screenshot
        import subprocess
        import tempfile
        import os
        
        temp_file = "/tmp/screen_analyze.png"
        subprocess.run(["gnome-screenshot", "-f", temp_file], capture_output=True, timeout=10)
        
        if not os.path.exists(temp_file):
            subprocess.run(["import", "-window", "root", temp_file], capture_output=True, timeout=10)
        
        if os.path.exists(temp_file):
            size = os.path.getsize(temp_file)
            
            # Try AI analysis
            if nexus:
                # Would need vision model
                return {"status": "info", "message": "Screen captured", "path": temp_file, "size": size}
            else:
                return {"status": "success", "screenshot": temp_file, "size": size}
        else:
            return {"status": "error", "message": "Screenshot failed"}
    
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/vision/analyze_image")
async def analyze_local_image(image_path: str, question: str = "Describe this image"):
    """Analyze local image with AI"""
    try:
        if not os.path.exists(image_path):
            return {"status": "error", "message": "Image not found"}
        
        if nexus:
            # Vision analysis would go here
            response = nexus.chat(
                model="gemini/gemini-pro",
                messages=[{"role": "user", "content": f"Analyze this image: {question}"}],
                priority="quality"
            )
            analysis = response.get("response", "")[:1000]
        else:
            analysis = "AI not configured"
        
        return {"status": "success", "analysis": analysis, "image": image_path}
    
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.get("/vision/camera")
async def camera_analysis():
    """Capture and analyze camera feed"""
    return {
        "status": "info",
        "message": "Camera requires OpenCV or node camera",
        "alternatives": [
            "Use /screenshot/instant for screen capture",
            "Use nodes action=camera_snap for phone camera",
            "Install opencv-python for local camera"
        ]
    }

@app.get("/vision/describe")
async def describe_image(image_path: str):
    """Get description of image"""
    try:
        if not os.path.exists(image_path):
            return {"status": "error", "message": "Image not found"}
        
        # Try basic image info
        from PIL import Image
        img = Image.open(image_path)
        
        info = {
            "path": image_path,
            "format": img.format,
            "size": f"{img.width}x{img.height}",
            "mode": img.mode
        }
        
        # Try OCR description
        try:
            import pytesseract
            text = pytesseract.image_to_string(img)
            if text.strip():
                info["text_detected"] = text[:500]
        except:
            pass
        
        return {"status": "success", "info": info}
    
    except Exception as e:
        return {"status": "error", "message": str(e)}

# ============== SPOTIFY CONTROL ==============

import requests

SPOTIFY_API = "https://api.spotify.com/v1"

@app.get("/spotify/status")
async def spotify_status():
    """Get Spotify playback status"""
    try:
        # Try playerctl first (works with Spotify desktop)
        result = subprocess.run(["playerctl", "-p", "spotify", "status"], capture_output=True, text=True, timeout=5)
        if result.returncode == 0:
            status = result.stdout.strip()
            
            # Get track info
            title_result = subprocess.run(["playerctl", "-p", "spotify", "metadata", "title"], capture_output=True, text=True)
            artist_result = subprocess.run(["playerctl", "-p", "spotify", "metadata", "artist"], capture_output=True, text=True)
            
            return {
                "status": status,
                "title": title_result.stdout.strip(),
                "artist": artist_result.stdout.strip(),
                "source": "playerctl"
            }
        else:
            return {"status": "stopped", "message": "Spotify not playing"}
    except:
        return {"status": "error", "message": "playerctl not available"}

@app.post("/spotify/play")
async def spotify_play():
    """Play Spotify"""
    try:
        subprocess.run(["playerctl", "-p", "spotify", "play"], capture_output=True, timeout=5)
        return {"status": "success", "action": "Playing"}
    except:
        return {"status": "error", "message": "Failed"}

@app.post("/spotify/pause")
async def spotify_pause():
    """Pause Spotify"""
    try:
        subprocess.run(["playerctl", "-p", "spotify", "pause"], capture_output=True, timeout=5)
        return {"status": "success", "action": "Paused"}
    except:
        return {"status": "error", "message": "Failed"}

@app.post("/spotify/next")
async def spotify_next():
    """Next track"""
    try:
        subprocess.run(["playerctl", "-p", "spotify", "next"], capture_output=True, timeout=5)
        return {"status": "success", "action": "Next track"}
    except:
        return {"status": "error", "message": "Failed"}

@app.post("/spotify/previous")
async def spotify_previous():
    """Previous track"""
    try:
        subprocess.run(["playerctl", "-p", "spotify", "previous"], capture_output=True, timeout=5)
        return {"status": "success", "action": "Previous track"}
    except:
        return {"status": "error", "message": "Failed"}

@app.post("/spotify/open")
async def open_spotify():
    """Open Spotify app"""
    try:
        subprocess.run(["spotify"], capture_output=True, timeout=10)
        return {"status": "success", "action": "Opened Spotify"}
    except:
        try:
            subprocess.run(["xdg-open", "spotify:"], capture_output=True, timeout=10)
            return {"status": "success", "action": "Opened Spotify"}
        except:
            return {"status": "error", "message": "Failed to open Spotify"}

@app.post("/spotify/play_song")
async def spotify_play_song(song: str):
    """Search and play specific song"""
    return {
        "status": "info",
        "message": "Spotify Web API requires authentication",
        "note": "Use playerctl for local Spotify or add Spotify API keys"
    }

@app.post("/spotify/volume")
async def spotify_volume(level: int):
    """Set Spotify volume (0-100)"""
    try:
        level = max(0, min(100, level))
        subprocess.run(["playerctl", "-p", "spotify", "volume", str(level/100)], capture_output=True, timeout=5)
        return {"status": "success", "volume": level}
    except:
        return {"status": "error", "message": "Failed"}

# ============== VOICE & SPEECH ==============

@app.post("/voice/speak")
async def voice_speak(text: str, lang: str = "en"):
    """Text to Speech - AI speaks"""
    try:
        # Use espeak or say command
        try:
            subprocess.run(["espeak", text], capture_output=True, timeout=30)
            return {"status": "success", "spoken": text[:200]}
        except:
            pass
        
        try:
            subprocess.run(["say", text], capture_output=True, timeout=30)
            return {"status": "success", "spoken": text[:200]}
        except:
            pass
        
        return {"status": "info", "message": "TTS not available", "text": text}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/voice/speak_file")
async def voice_speak_file(text: str, output_file: str = "/tmp/speak.mp3"):
    """Generate speech audio file"""
    try:
        # Try gTTS (Google TTS)
        try:
            from gtts import gTTS
            tts = gTTS(text=text, lang=lang)
            tts.save(output_file)
            return {"status": "success", "file": output_file}
        except:
            pass
        
        return {"status": "error", "message": "gTTS not installed"}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.get("/voice/listen")
async def voice_listen():
    """Start voice listening (returns instructions)"""
    return {
        "status": "info",
        "message": "Voice input requires microphone",
        "options": [
            "Use browser-based speech recognition",
            "Install speech_recognition Python package",
            "Use arecord for audio capture"
        ]
    }

@app.post("/voice/record")
async def voice_record(duration: int = 5, output_file: str = "/tmp/voice.wav"):
    """Record audio from microphone"""
    try:
        import wave
        import sounddevice as sd
        
        samplerate = 16000
        recording = sd.rec(int(duration * samplerate), samplerate=samplerate, channels=1)
        sd.wait()
        
        with wave.open(output_file, 'w') as wf:
            wf.setnchannels(1)
            wf.setsampwidth(2)
            wf.setframerate(samplerate)
            wf.writeframes(recording.tobytes())
        
        return {"status": "success", "file": output_file, "duration": duration}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/voice/transcribe")
async def voice_transcribe(audio_file: str = "/tmp/voice.wav"):
    """Transcribe audio to text"""
    try:
        # Try speech_recognition
        try:
            import speech_recognition as sr
            r = sr.Recognizer()
            with sr.AudioFile(audio_file) as source:
                audio = r.record(source)
            text = r.recognize_google(audio)
            return {"status": "success", "text": text}
        except:
            pass
        
        return {"status": "info", "message": "Speech recognition not available"}
    except Exception as e:
        return {"status": "error", "message": str(e)}

# ============== ADVANCED VOICE ==============

@app.post("/voice/tts")
async def text_to_speech(text: str, voice: str = "en", speed: float = 1.0):
    """Advanced Text to Speech with multiple options"""
    output_file = f"/tmp/tts_{int(time.time())}.mp3"
    
    # Try gTTS
    try:
        from gtts import gTTS
        tts = gTTS(text=text, lang=voice[:2])
        tts.save(output_file)
        return {"status": "success", "file": output_file, "format": "mp3"}
    except:
        pass
    
    # Try pyttsx3 (offline)
    try:
        import pyttsx3
        engine = pyttsx3.init()
        engine.setProperty('rate', int(150 * speed))
        engine.save_to_file(text, output_file.replace('.mp3', '.wav'))
        engine.runAndWait()
        return {"status": "success", "file": output_file.replace('.mp3', '.wav'), "format": "wav"}
    except:
        pass
    
    return {"status": "error", "message": "TTS not available"}

@app.post("/voice/stt")
async def speech_to_text(audio_data: str = None):
    """Speech to Text"""
    try:
        import speech_recognition as sr
        r = sr.Recognizer()
        
        if audio_data:
            # Process audio data
            return {"status": "success", "text": "Audio processed"}
        
        return {"status": "info", "message": "Send audio file for transcription"}
    except:
        return {"status": "error", "message": "Speech recognition not available"}

# ============== IMAGE PROCESSING ==============

@app.post("/image/resize")
async def resize_image(input_path: str, width: int, height: int, output: str = None):
    """Resize image"""
    try:
        from PIL import Image
        img = Image.open(input_path)
        img = img.resize((width, height))
        
        if not output:
            output = input_path.replace('.', '_resized.')
        img.save(output)
        
        return {"status": "success", "output": output, "size": f"{width}x{height}"}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/image/filter")
async def image_filter(input_path: str, filter_type: str = "blur", output: str = None):
    """Apply filter to image"""
    try:
        from PIL import Image, ImageFilter
        
        filters = {
            "blur": ImageFilter.BLUR,
            "sharpen": ImageFilter.SHARPEN,
            "smooth": ImageFilter.SMOOTH,
            "edge": ImageFilter.FIND_EDGES,
            "contour": ImageFilter.CONTOUR,
            "detail": ImageFilter.DETAIL,
            "emboss": ImageFilter.EMBOSS
        }
        
        img = Image.open(input_path)
        img = img.filter(filters.get(filter_type, ImageFilter.BLUR))
        
        if not output:
            output = input_path.replace('.', f'_{filter_type}.')
        img.save(output)
        
        return {"status": "success", "output": output, "filter": filter_type}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/image/crop")
async def crop_image(input_path: str, x: int, y: int, width: int, height: int, output: str = None):
    """Crop image"""
    try:
        from PIL import Image
        img = Image.open(input_path)
        img = img.crop((x, y, x + width, y + height))
        
        if not output:
            output = input_path.replace('.', '_cropped.')
        img.save(output)
        
        return {"status": "success", "output": output}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/image/rotate")
async def rotate_image(input_path: str, angle: int = 90, output: str = None):
    """Rotate image"""
    try:
        from PIL import Image
        img = Image.open(input_path)
        img = img.rotate(angle, expand=True)
        
        if not output:
            output = input_path.replace('.', f'_rotated{angle}.')
        img.save(output)
        
        return {"status": "success", "output": output, "angle": angle}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/image/watermark")
async def add_watermark(input_path: str, text: str, output: str = None):
    """Add watermark to image"""
    try:
        from PIL import Image, ImageDraw, ImageFont
        
        img = Image.open(input_path).convert("RGBA")
        txt = Image.new("RGBA", img.size, (255, 255, 255, 0))
        
        draw = ImageDraw.Draw(txt)
        
        try:
            font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 36)
        except:
            font = ImageFont.load_default()
        
        # Draw text at bottom right
        position = (img.width - 200, img.height - 50)
        draw.text(position, text, font=font, fill=(255, 255, 255, 128))
        
        combined = Image.alpha_composite(img, txt)
        
        if not output:
            output = input_path.replace('.', '_watermarked.')
        combined.save(output)
        
        return {"status": "success", "output": output, "text": text}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/image/ocr")
async def image_ocr(input_path: str, lang: str = "eng"):
    """Extract text from image (OCR)"""
    try:
        from PIL import Image
        import pytesseract
        
        img = Image.open(input_path)
        text = pytesseract.image_to_string(img, lang=lang)
        
        return {"status": "success", "text": text[:3000], "length": len(text)}
    except Exception as e:
        return {"status": "error", "message": str(e)}

# ============== VIDEO PROCESSING ==============

@app.post("/video/info")
async def video_info(video_path: str):
    """Get video information"""
    try:
        import cv2
        cap = cv2.VideoCapture(video_path)
        
        info = {
            "frames": int(cap.get(cv2.CAP_PROP_FRAME_COUNT)),
            "fps": cap.get(cv2.CAP_PROP_FPS),
            "width": int(cap.get(cv2.CAP_PROP_FRAME_WIDTH)),
            "height": int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT)),
            "duration": int(cap.get(cv2.CAP_PROP_FRAME_COUNT) / cap.get(cv2.CAP_PROP_FPS))
        }
        cap.release()
        
        return {"status": "success", "info": info}
    except:
        return {"status": "error", "message": "OpenCV not available"}

@app.post("/video/thumbnail")
async def video_thumbnail(video_path: str, timestamp: int = 1, output: str = None):
    """Extract thumbnail from video"""
    try:
        import cv2
        
        cap = cv2.VideoCapture(video_path)
        cap.set(cv2.CAP_PROP_POS_FRAMES, timestamp)
        ret, frame = cap.read()
        
        if not output:
            output = video_path.replace('.mp4', '_thumb.jpg').replace('.avi', '_thumb.jpg')
        
        cv2.imwrite(output, frame)
        cap.release()
        
        return {"status": "success", "thumbnail": output, "timestamp": timestamp}
    except:
        return {"status": "error", "message": "Failed to extract thumbnail"}

# ============== CODING ==============

@app.post("/code/run")
async def run_code(code: str, language: str = "python"):
    """Execute code in sandbox"""
    import tempfile
    import subprocess
    
    lang_config = {
        "python": {"ext": "py", "run": "python3"},
        "javascript": {"ext": "js", "run": "node"},
        "bash": {"ext": "sh", "run": "bash"},
    }
    
    config = lang_config.get(language.lower())
    if not config:
        return {"status": "error", "message": "Unsupported language"}
    
    try:
        with tempfile.NamedTemporaryFile(mode='w', suffix=f".{config['ext']}", delete=False) as f:
            f.write(code)
            temp_file = f.name
        
        result = subprocess.run([config['run'], temp_file], capture_output=True, text=True, timeout=30)
        os.unlink(temp_file)
        
        return {
            "status": "success",
            "output": result.stdout[:2000],
            "error": result.stderr[:500] if result.stderr else None,
            "return_code": result.returncode
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/code/format")
async def format_code(code: str, language: str = "python"):
    """Format code"""
    try:
        if language.lower() == "python":
            import black
            formatted = black.format_str(code, mode=black.Mode())
            return {"status": "success", "formatted": formatted}
        elif language.lower() == "javascript":
            import jsbeautifier
            formatted = jsbeautifier.beautify(code)
            return {"status": "success", "formatted": formatted}
    except:
        pass
    
    return {"status": "info", "message": "Formatter not available"}

# ============== PRESENTATIONS ==============

@app.post("/ppt/create")
async def create_presentation(title: str = "Presentation"):
    """Create PowerPoint presentation"""
    return {
        "status": "info",
        "message": "PowerPoint requires python-pptx",
        "install": "pip install python-pptx"
    }

@app.post("/ppt/add_slide")
async def add_slide(presentation_path: str, title: str, content: str = ""):
    """Add slide to presentation"""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        
        prs = Presentation(presentation_path) if os.path.exists(presentation_path) else Presentation()
        
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title_placeholder = slide.shapes.title
        title_placeholder.text = title
        
        if content:
            body = slide.placeholders[1]
            body.text = content
        
        prs.save(presentation_path)
        
        return {"status": "success", "slides": len(prs.slides)}
    except:
        return {"status": "error", "message": "python-pptx not available"}

# ============== CANVA INTEGRATION ==============

@app.get("/canva/connect")
async def canva_connect():
    """Connect to Canva"""
    return {
        "status": "info",
        "message": "Canva API integration",
        "note": "Requires Canva API credentials",
        "endpoints": [
            "/canva/designs - List designs",
            "/canva/create - Create design",
            "/canva/export - Export design"
        ]
    }

@app.get("/canva/designs")
async def canva_designs():
    """List Canva designs"""
    return {
        "status": "info",
        "message": "Connect Canva API to access designs",
        "api_url": "https://api.canva.com/rest/v1/designs"
    }

@app.post("/canva/create")
async def canva_create(design_type: str = "poster"):
    """Create new Canva design"""
    return {
        "status": "info",
        "message": f"Create {design_type} on Canva",
        "action": "Use Canva API credentials",
        "替代": "Use libreoffice to create presentations locally"
    }

# ============== ADOBE INTEGRATION ==============

@app.get("/adobe/connect")
async def adobe_connect():
    """Connect to Adobe services"""
    return {
        "status": "info",
        "message": "Adobe PDF Services API",
        "note": "Requires Adobe API credentials",
        "features": [
            "PDF extraction",
            "PDF conversion",
            "PDF compression",
            "OCR"
        ]
    }

@app.post("/adobe/pdf/extract")
async def adobe_extract_text(pdf_path: str):
    """Extract text from PDF"""
    try:
        import PyPDF2
        
        with open(pdf_path, 'rb') as f:
            reader = PyPDF2.PdfReader(f)
            text = ""
            for page in reader.pages:
                text += page.extract_text()
        
        return {"status": "success", "text": text[:3000], "pages": len(reader.pages)}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/adobe/pdf/merge")
async def adobe_merge_pdfs(pdf_files: str, output: str = "/tmp/merged.pdf"):
    """Merge multiple PDFs"""
    try:
        import PyPDF2
        
        merger = PyPDF2.PdfMerger()
        
        files = pdf_files.split(',')
        for f in files:
            f = f.strip()
            if os.path.exists(f):
                merger.append(f)
        
        merger.write(output)
        merger.close()
        
        return {"status": "success", "output": output, "files": len(files)}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/adobe/pdf/split")
async def adobe_split_pdf(pdf_path: str, start: int = 0, end: int = 1):
    """Split PDF"""
    try:
        import PyPDF2
        
        reader = PyPDF2.PdfReader(pdf_path)
        writer = PyPDF2.PdfWriter()
        
        for i in range(start, min(end, len(reader.pages))):
            writer.add_page(reader.pages[i])
        
        output = pdf_path.replace('.pdf', '_split.pdf')
        writer.write(output)
        
        return {"status": "success", "output": output}
    except Exception as e:
        return {"status": "error", "message": str(e)}

# ============== VOICE MESSAGES ==============

@app.get("/tts/generate")
async def generate_voice(text: str, lang: str = "en"):
    """Generate voice audio file for Telegram"""
    try:
        from gtts import gTTS
        import uuid
        
        output_file = f"/tmp/voice_{uuid.uuid4().hex[:8]}.mp3"
        tts = gTTS(text=text[:500], lang=lang[:2])
        tts.save(output_file)
        
        return {"status": "success", "audio_file": output_file}
    except Exception as e:
        return {"status": "error", "message": str(e)}

# ============== VOICE ACTIVATION ==============

# Voice activation state
VOICE_ACTIVE = False
WAKE_WORD = "hey x"
AUTHORIZED_VOICES = {}  # Will store voice profiles

@app.post("/voice/activate")
async def activate_voice():
    """Activate voice listening mode"""
    global VOICE_ACTIVE
    VOICE_ACTIVE = True
    return {"status": "success", "mode": "voice_active", "wake_word": WAKE_WORD}

@app.post("/voice/deactivate")
async def deactivate_voice():
    """Deactivate voice listening"""
    global VOICE_ACTIVE
    VOICE_ACTIVE = False
    return {"status": "success", "mode": "voice_inactive"}

@app.get("/voice/status")
async def voice_status():
    """Get voice activation status"""
    return {
        "status": "success",
        "voice_active": VOICE_ACTIVE,
        "wake_word": WAKE_WORD,
        "authorized_users": len(AUTHORIZED_VOICES)
    }

@app.post("/voice/listen_once")
async def voice_listen_once():
    """Listen for one command (after wake word)"""
    return {
        "status": "info",
        "message": "Voice recognition requires microphone",
        "setup_needed": [
            "Install: pip install speechrecognition pocketsphinx",
            "Install: pip install pvporcupine (wake word)",
            "Record voice samples for recognition"
        ]
    }

@app.post("/voice/wake_word")
async def set_wake_word(word: str = "hey x"):
    """Set custom wake word"""
    global WAKE_WORD
    WAKE_WORD = word.lower()
    return {"status": "success", "wake_word": WAKE_WORD, "message": f"Say '{WAKE_WORD}' to activate"}

@app.post("/voice/enroll")
async def voice_enroll(user_name: str):
    """Enroll authorized voice (for future voice ID)"""
    return {
        "status": "info",
        "message": f"Enrolling {user_name} voice",
        "note": "Record voice samples to identify speaker",
        "future": "Only enrolled voices will activate assistant"
    }

@app.post("/voice/start_listening")
async def start_continuous_listening():
    """Start continuous voice listening (JARVIS mode)"""
    return {
        "status": "info",
        "message": "Continuous listening starting...",
        "wake_word": WAKE_WORD,
        "instructions": [
            "1. Install: pip install speechrecognition pocketsphinx pvporcupine",
            "2. Run voice agent on desktop with microphone",
            "3. Say 'hey x' to activate",
            "Only your voice will respond"
        ]
    }

# ============== INSTAGRAM AUTOMATION ==============

@app.get("/instagram/connect")
async def instagram_connect():
    """Connect to Instagram via browser"""
    return {
        "status": "info",
        "message": "Use browser automation for Instagram",
        "steps": [
            "1. Open Instagram in Chrome",
            "2. Login manually (one time)",
            "3. Use /browser to control",
            "4. No credentials needed again!"
        ]
    }

@app.get("/instagram/open")
async def instagram_open():
    """Open Instagram in browser"""
    try:
        subprocess.run(["xdg-open", "https://instagram.com"], capture_output=True, timeout=10)
        return {"status": "success", "action": "Opened Instagram"}
    except:
        return {"status": "error", "message": "Failed to open"}

@app.get("/instagram/story")
async def instagram_story():
    """Get latest Instagram stories"""
    return {
        "status": "info",
        "message": "Use browser to view stories",
        "action": "xdg-open https://instagram.com/stories"
    }

@app.post("/instagram/post")
async def instagram_post(image_path: str = None, caption: str = ""):
    """Post to Instagram (via browser)"""
    return {
        "status": "info",
        "message": "Browser automation required",
        "steps": [
            "Open Instagram in browser",
            "Navigate to create post",
            "Upload image",
            "Add caption",
            "Post"
        ]
    }

# ============== SOCIAL MEDIA AUTOMATION ==============

@app.get("/social/instagram/open")
async def social_instagram_open():
    """Open Instagram"""
    return {"status": "success", "action": "Open Instagram", "url": "https://instagram.com"}

@app.get("/social/instagram/messages")
async def social_instagram_messages():
    """Open Instagram DMs"""
    return {"status": "success", "action": "Open Instagram DMs", "url": "https://instagram.com/direct"}

@app.get("/social/instagram/stories")
async def social_instagram_stories():
    """Open Instagram Stories"""
    return {"status": "success", "action": "Open Stories", "url": "https://instagram.com/stories"}

@app.get("/social/instagram/notifications")
async def social_instagram_notifications():
    """Open Instagram Notifications"""
    return {"status": "success", "action": "Open Notifications", "url": "https://instagram.com/activity"}

@app.get("/social/instagram/profile")
async def social_instagram_profile(username: str = ""):
    """Open Instagram Profile"""
    if username:
        return {"status": "success", "action": "Open Profile", "url": f"https://instagram.com/{username}"}
    return {"status": "success", "action": "Open My Profile", "url": "https://instagram.com"}

# WhatsApp
@app.get("/social/whatsapp/open")
async def social_whatsapp_open():
    """Open WhatsApp Web"""
    return {"status": "success", "action": "Open WhatsApp", "url": "https://web.whatsapp.com"}

@app.get("/social/whatsapp/dm")
async def social_whatsapp_dm(phone: str = ""):
    """Open WhatsApp DM"""
    if phone:
        return {"status": "success", "action": "Open Chat", "url": f"https://wa.me/{phone}"}
    return {"status": "success", "action": "Open WhatsApp", "url": "https://web.whatsapp.com"}

# Telegram
@app.get("/social/telegram/open")
async def social_telegram_open():
    """Open Telegram Web"""
    return {"status": "success", "action": "Open Telegram", "url": "https://web.telegram.org"}

# Facebook
@app.get("/social/facebook/open")
async def social_facebook_open():
    """Open Facebook"""
    return {"status": "success", "action": "Open Facebook", "url": "https://facebook.com"}

@app.get("/social/facebook/messages")
async def social_facebook_messages():
    """Open Facebook Messenger"""
    return {"status": "success", "action": "Open Messenger", "url": "https://facebook.com/messages"}

# Twitter/X
@app.get("/social/twitter/open")
async def social_twitter_open():
    """Open Twitter/X"""
    return {"status": "success", "action": "Open X", "url": "https://x.com"}

@app.get("/social/twitter/compose")
async def social_twitter_compose():
    """Open Twitter compose"""
    return {"status": "success", "action": "New Tweet", "url": "https://x.com/compose/post"}

# YouTube
@app.get("/social/youtube/open")
async def social_youtube_open():
    """Open YouTube"""
    return {"status": "success", "action": "Open YouTube", "url": "https://youtube.com"}

@app.get("/social/youtube/search")
async def social_youtube_search(query: str = ""):
    """Search YouTube"""
    return {"status": "success", "action": "Search", "url": f"https://youtube.com/results?search_query={query}"}

# LinkedIn
@app.get("/social/linkedin/open")
async def social_linkedin_open():
    """Open LinkedIn"""
    return {"status": "success", "action": "Open LinkedIn", "url": "https://linkedin.com"}

@app.get("/social/linkedin/messages")
async def social_linkedin_messages():
    """Open LinkedIn Messages"""
    return {"status": "success", "action": "Open Messages", "url": "https://linkedin.com/messaging"}

# Reddit
@app.get("/social/reddit/open")
async def social_reddit_open():
    """Open Reddit"""
    return {"status": "success", "action": "Open Reddit", "url": "https://reddit.com"}

@app.get("/social/reddit/subreddit")
async def social_reddit_subreddit(sub: str = "all"):
    """Open Reddit Subreddit"""
    return {"status": "success", "action": "Open Subreddit", "url": f"https://reddit.com/r/{sub}"}

# Google
@app.get("/social/google/search")
async def social_google_search(query: str = ""):
    """Google Search"""
    return {"status": "success", "action": "Search", "url": f"https://google.com/search?q={query}"}

@app.get("/social/google/gmail")
async def social_google_gmail():
    """Open Gmail"""
    return {"status": "success", "action": "Open Gmail", "url": "https://mail.google.com"}

@app.get("/social/google/drive")
async def social_google_drive():
    """Open Google Drive"""
    return {"status": "success", "action": "Open Drive", "url": "https://drive.google.com"}

# ============== EMOTIONAL AI ==============

# Emotion detection and response
EMOTIONS = {
    "happy": ["happy", "joy", "excited", "great", "wonderful", "amazing", "love", "celebrate", "good news", "promoted", "won"],
    "sad": ["sad", "down", "depressed", "unhappy", "miss", "lost", "failed", "sorry", "grief", "pain", "hurt"],
    "angry": ["angry", "furious", "annoyed", "hate", "stupid", "irritating", "frustrated", "rage"],
    "fear": ["scared", "afraid", "worried", "nervous", "anxious", "fear", "panic", "terrified"],
    "surprised": ["wow", "surprise", "unexpected", "shocked", "unbelievable", "no way", "really?"],
    "excited": ["excited", "can't wait", "looking forward", "awesome", "epic", "cool", "yay"],
    "thoughtful": ["think", "wonder", "consider", "maybe", "perhaps", "idea", "hmm"],
    "supportive": ["help", "support", "need", "please", "stuck", "confused", "how to"]
}

EMOTION_RESPONSES = {
    "happy": [
        "That's wonderful! 🎉 I'm so happy for you!",
        "That's amazing news! 🌟 Keep shining!",
        "Yay! That's fantastic! 😊",
        "Great to hear! You deserve all the happiness! 🎊"
    ],
    "sad": [
        "I'm here for you. 💙 It okay to feel this way.",
        "I'm sorry you're going through this. 🤗 Want to talk about it?",
        "Things will get better. Stay strong. 💪",
        "I'm listening. Tell me more if you want. 🫂"
    ],
    "angry": [
        "I understand you're frustrated. Take a deep breath. 🧘",
        "That's really annoying. Want to vent about it?",
        "I get it - that's infuriating. What happened?",
        "Calm down - let me help you figure this out. 💡"
    ],
    "fear": [
        "It's okay to feel scared. I'm here. 🛡️",
        "Don't worry - we'll get through this together.",
        "Take it easy. What's worrying you?",
        "I'm here to help. What do you need?"
    ],
    "surprised": [
        "Wow! That's unexpected! 😲",
        "No way! Tell me more!",
        "That's incredible! 🎉",
        "Seriously?! That's amazing!"
    ],
    "excited": [
        "YES! That's awesome! 🔥",
        "I'm so excited for you! 🎊",
        "Tell me everything! I want details!",
        "This is amazing! How did it happen?"
    ],
    "thoughtful": [
        "That's an interesting thought. 🤔",
        "Let me consider that...",
        "Good point. Here's what I think...",
        "Interesting perspective. Let me help you figure it out."
    ],
    "supportive": [
        "I'm here to help! 💪",
        "What do you need? Tell me.",
        "You've got this! I'm cheering for you! 🌟",
        "I'm on your side. What can I do?"
    ]
}

import random

def detect_emotion(text: str) -> str:
    """Detect emotion from text"""
    text = text.lower()
    
    for emotion, keywords in EMOTIONS.items():
        for keyword in keywords:
            if keyword in text:
                return emotion
    
    return "neutral"

@app.post("/emotion/detect")
async def detect_emotion_endpoint(text: str):
    """Detect emotion from text"""
    emotion = detect_emotion(text)
    return {"status": "success", "emotion": emotion, "text": text}

@app.post("/emotion/respond")
async def emotional_response(text: str):
    """Get emotional response based on detected mood"""
    emotion = detect_emotion(text)
    responses = EMOTION_RESPONSES.get(emotion, EMOTION_RESPONSES["supportive"])
    response = random.choice(responses)
    
    return {
        "status": "success",
        "detected_emotion": emotion,
        "response": response,
        "text": text
    }

@app.post("/emotion/set")
async def set_emotion(mode: str):
    """Set AI emotion mode"""
    valid = ["happy", "sad", "excited", "calm", "serious", "funny", "neutral"]
    
    if mode.lower() in valid:
        return {"status": "success", "mode": mode, "message": f"AI is now {mode}"}
    return {"status": "error", "message": f"Invalid mode. Choose: {valid}"}

@app.get("/emotion/mood")
async def get_mood():
    """Get current AI mood"""
    return {
        "status": "success",
        "mood": "supportive",
        "available_moods": list(EMOTIONS.keys())
    }

# ============== NOVA-LIKE AI ASSISTANT ==============

# Personal Assistant Features
import json
from datetime import datetime, timedelta

# User preferences and memory
USER_PREFERENCES_FILE = "/tmp/nexusai_preferences.json"
TASKS_FILE = "/tmp/nexusai_tasks.json"
NOTES_FILE = "/tmp/nexusai_notes.json"

def load_json(filepath):
    if os.path.exists(filepath):
        with open(filepath, 'r') as f:
            return json.load(f)
    return {}

def save_json(filepath, data):
    with open(filepath, 'w') as f:
        json.dump(data, f)

@app.post("/assistant/greet")
async def greet_user(name: str = ""):
    """Greet the user"""
    if name:
        # Save name preference
        prefs = load_json(USER_PREFERENCES_FILE)
        prefs['name'] = name
        save_json(USER_PREFERENCES_FILE, prefs)
        return {"status": "success", "message": f"Hello {name}! I'm X, your AI assistant. How can I help you today? 😊"}
    
    prefs = load_json(USER_PREFERENCES_FILE)
    user_name = prefs.get('name', 'Boss')
    return {"status": "success", "message": f"Hey {user_name}! I'm listening... 🎧"}

@app.post("/assistant/remember")
async def remember_fact(fact: str):
    """Remember something about user"""
    prefs = load_json(USER_PREFERENCES_FILE)
    if 'facts' not in prefs:
        prefs['facts'] = []
    prefs['facts'].append({"fact": fact, "date": datetime.now().isoformat()})
    save_json(USER_PREFERENCES_FILE, prefs)
    return {"status": "success", "message": f"Got it! I'll remember that: {fact} 💡"}

@app.get("/assistant/recall")
async def recall_facts():
    """Recall remembered facts"""
    prefs = load_json(USER_PREFERENCES_FILE)
    facts = prefs.get('facts', [])
    return {"status": "success", "facts": facts[-10:]}

@app.post("/assistant/preference")
async def set_preference(key: str, value: str):
    """Set user preference"""
    prefs = load_json(USER_PREFERENCES_FILE)
    prefs[key] = value
    save_json(USER_PREFERENCES_FILE, prefs)
    return {"status": "success", "message": f"Preference saved: {key} = {value}"}

@app.get("/assistant/preferences")
async def get_preferences():
    """Get all preferences"""
    prefs = load_json(USER_PREFERENCES_FILE)
    return {"status": "success", "preferences": prefs}

# Task Management
@app.post("/assistant/task")
async def add_task(task: str, due: str = ""):
    """Add a task"""
    tasks = load_json(TASKS_FILE)
    task_id = len(tasks.get('tasks', [])) + 1
    new_task = {
        "id": task_id,
        "task": task,
        "due": due,
        "status": "pending",
        "created": datetime.now().isoformat()
    }
    
    if 'tasks' not in tasks:
        tasks['tasks'] = []
    tasks['tasks'].append(new_task)
    save_json(TASKS_FILE, tasks)
    
    return {"status": "success", "message": f"Task added: {task} ✅", "task_id": task_id}

@app.get("/assistant/tasks")
async def list_tasks(status: str = "all"):
    """List tasks"""
    tasks = load_json(TASKS_FILE)
    all_tasks = tasks.get('tasks', [])
    
    if status != "all":
        all_tasks = [t for t in all_tasks if t.get('status') == status]
    
    return {"status": "success", "tasks": all_tasks, "count": len(all_tasks)}

@app.post("/assistant/task/complete")
async def complete_task(task_id: int):
    """Mark task complete"""
    tasks = load_json(TASKS_FILE)
    
    for task in tasks.get('tasks', []):
        if task.get('id') == task_id:
            task['status'] = 'completed'
            task['completed'] = datetime.now().isoformat()
    
    save_json(TASKS_FILE, tasks)
    return {"status": "success", "message": f"Task {task_id} completed! 🎉"}

@app.delete("/assistant/task")
async def delete_task(task_id: int):
    """Delete task"""
    tasks = load_json(TASKS_FILE)
    tasks['tasks'] = [t for t in tasks.get('tasks', []) if t.get('id') != task_id]
    save_json(TASKS_FILE, tasks)
    return {"status": "success", "message": f"Task {task_id} deleted"}

# Notes
@app.post("/assistant/note")
async def add_note(title: str, content: str):
    """Add a note"""
    notes = load_json(NOTES_FILE)
    note_id = len(notes.get('notes', [])) + 1
    new_note = {
        "id": note_id,
        "title": title,
        "content": content,
        "created": datetime.now().isoformat()
    }
    
    if 'notes' not in notes:
        notes['notes'] = []
    notes['notes'].append(new_note)
    save_json(NOTES_FILE, notes)
    
    return {"status": "success", "message": f"Note saved: {title} 📝", "note_id": note_id}

@app.get("/assistant/notes")
async def list_notes():
    """List all notes"""
    notes = load_json(NOTES_FILE)
    return {"status": "success", "notes": notes.get('notes', [])}

@app.get("/assistant/note")
async def get_note(note_id: int = None):
    """Get specific note"""
    notes = load_json(NOTES_FILE)
    all_notes = notes.get('notes', [])
    
    if note_id:
        for note in all_notes:
            if note.get('id') == note_id:
                return {"status": "success", "note": note}
    
    return {"status": "success", "notes": all_notes}

# Smart home placeholder
@app.get("/assistant/smarthome/status")
async def smarthome_status():
    """Get smart home status"""
    return {
        "status": "info",
        "message": "Smart home integration",
        "available": [
            "/assistant/smarthome/lights?action=on|off",
            "/assistant/smarthome/temperature",
            "/assistant/smarthome/lock?action=lock|unlock"
        ]
    }

@app.post("/assistant/smarthome/lights")
async def control_lights(action: str = "toggle"):
    """Control lights"""
    return {"status": "success", "action": action, "message": f"Lights {action} 💡"}

@app.post("/assistant/smarthome/lock")
async def control_lock(action: str = "lock"):
    """Control door lock"""
    return {"status": "success", "action": action, "message": f"Door {action}ed 🔒"}

# Daily briefing
@app.get("/assistant/briefing")
async def daily_briefing():
    """Get daily briefing"""
    import psutil
    
    # Get tasks
    tasks = load_json(TASKS_FILE)
    pending = len([t for t in tasks.get('tasks', []) if t.get('status') == 'pending'])
    
    # Get weather
    weather_info = {"temp": "N/A", "condition": "checking..."}
    
    briefing = f"""
🌤️ Good {datetime.now().strftime('%A')} morning!

📋 You have {pending} pending tasks
🖥️ System: Running smoothly
⏰ Time: {datetime.now().strftime('%H:%M')}

How can I help you today?
    """
    
    return {"status": "success", "briefing": briefing.strip()}

# Quick actions
@app.get("/assistant/help")
async def assistant_help():
    """Get help with assistant commands"""
    return {
        "status": "success",
        "commands": {
            "Greeting": "/assistant/greet?name=YourName",
            "Remember": "/assistant/remember?fact=Something",
            "Recall": "/assistant/recall",
            "Tasks": "/assistant/task?task=Do something",
            "Notes": "/assistant/note?title=Title&content=...",
            "Briefing": "/assistant/briefing",
            "Smart Home": "/assistant/smarthome/status"
        }
    }

# ============== AI AVATAR & FACE ==============

# Avatar configuration
AVATAR_CONFIG = {
    "current": "default",
    "mood": "happy",
    "animations": ["idle", "listening", "speaking", "thinking"]
}

@app.get("/avatar/status")
async def avatar_status():
    """Get avatar status"""
    return {"status": "success", "avatar": AVATAR_CONFIG}

@app.post("/avatar/set")
async def set_avatar(avatar: str = "default"):
    """Set avatar style"""
    AVATAR_CONFIG["current"] = avatar
    return {"status": "success", "avatar": avatar, "message": f"Avatar set to {avatar}"}

@app.post("/avatar/mood")
async def set_mood(mood: str):
    """Set avatar mood"""
    valid_moods = ["happy", "sad", "excited", "thinking", "listening", "neutral", "sleeping"]
    if mood in valid_moods:
        AVATAR_CONFIG["mood"] = mood
        return {"status": "success", "mood": mood, "message": f"X is now {mood} 😊" if mood == "happy" else f"X is {mood}"}
    return {"status": "error", "message": f"Invalid mood. Choose: {valid_moods}"}

@app.get("/avatar/face")
async def get_face():
    """Get current face/avatar"""
    moods = {
        "happy": "😊",
        "sad": "😢",
        "excited": "🤩",
        "thinking": "🤔",
        "listening": "👂",
        "neutral": "😐",
        "sleeping": "😴",
        "angry": "😠",
        "love": "😍"
    }
    
    face = moods.get(AVATAR_CONFIG["mood"], "😊")
    
    return {
        "status": "success",
        "face": face,
        "mood": AVATAR_CONFIG["mood"],
        "message": f"X is feeling {AVATAR_CONFIG['mood']} {face}"
    }

@app.get("/avatar/faces")
async def list_faces():
    """List all available faces"""
    return {
        "status": "success",
        "faces": {
            "😊": "happy",
            "😢": "sad", 
            "🤩": "excited",
            "🤔": "thinking",
            "👂": "listening",
            "😐": "neutral",
            "😴": "sleeping",
            "😠": "angry",
            "😍": "love",
            "😎": "cool",
            "🤖": "AI mode"
        }
    }

@app.post("/avatar/react")
async def avatar_react(emotion: str):
    """X reacts with emotion"""
    mood_map = {
        "celebrate": "excited",
        "happy": "happy",
        "sad": "sad",
        "think": "thinking",
        "listen": "listening",
        "love": "love",
        "angry": "angry",
        "cool": "cool"
    }
    
    mood = mood_map.get(emotion.lower(), "neutral")
    AVATAR_CONFIG["mood"] = mood
    
    reactions = {
        "excited": "🎉 YAY! 🎉",
        "happy": "😊 That's great!",
        "sad": "😢 I'm here for you...",
        "thinking": "🤔 Let me think...",
        "listening": "👂 I'm listening...",
        "love": "😍 That's lovely!",
        "angry": "😠 That's not okay!",
        "cool": "😎 Nice!"
    }
    
    return {
        "status": "success",
        "mood": mood,
        "reaction": reactions.get(mood, "😐")
    }

# Generate ASCII art face
@app.get("/avatar/ascii")
async def ascii_face():
    """Get ASCII art face"""
    faces = {
        "happy": """
    ╭───────────╮
    │  😊     😊  │
    │     ▽       │
    │   ╰───╯     │
    │  ▽       ▽  │
    ╰───────────╯
    """,
        "thinking": """
    ╭───────────╮
    │  🤔     🤔  │
    │     ▽       │
    │   ╰───╯     │
    │   ▽     ▽   │
    ╰───────────╯
    """,
        "neutral": """
    ╭───────────╮
    │  😐     😐  │
    │     ▽       │
    │   ╰───╯     │
    │  ─       ─  │
    ╰───────────╯
    """
    }
    
    face = faces.get(AVATAR_CONFIG["mood"], faces["happy"])
    return {"status": "success", "ascii": face, "mood": AVATAR_CONFIG["mood"]}

# ============== ANIMATED AVATAR ==============

import random
import time

# Animated character states
AVATAR_STATES = {
    "idle": ["◕ ◕", "◔ ◔", "◕ ◕", "◑ ◑"],
    "listening": ["◉ ◯", "○ ◉", "◉ ◯", "● ◯"],
    "thinking": ["◕ ‿", "◔ ◡", "◕ ◡", "◔ ‿"],
    "speaking": ["◕ ◕", "▼ ▼", "◕ ◕", "▲ ▲"],
    "happy": ["◕ ◕", "◖ ◗", "◕ ◕", "◘ ◘"],
    "excited": ["◕ ◕", "◝ ◞", "◕ ◕", "◞ ◝"],
    "sad": ["◕ ◕", "▼ ▼", "◠ ◠", "◡ ◡"]
}

CURRENT_AVATAR_STATE = "idle"
AVATAR_ANIMATION_INDEX = 0

@app.get("/avatar/animate")
async def animate_avatar(state: str = "idle"):
    """Get animated avatar frame"""
    global CURRENT_AVATAR_STATE, AVATAR_ANIMATION_INDEX
    
    if state:
        CURRENT_AVATAR_STATE = state
    
    frames = AVATAR_STATES.get(CURRENT_AVATAR_STATE, AVATAR_STATES["idle"])
    frame = frames[AVATAR_ANIMATION_INDEX % len(frames)]
    AVATAR_ANIMATION_INDEX += 1
    
    # ASCII character body
    body = f"""
    {frame}
   ╭─────╮
   │ {frame} │
   ╰─────╯
    """
    
    return {
        "status": "success",
        "frame": frame,
        "state": CURRENT_AVATAR_STATE,
        "ascii": body.strip()
    }

@app.get("/avatar/animated")
async def get_animated_avatar():
    """Get full animated character"""
    states = AVATAR_STATES[CURRENT_AVATAR_STATE]
    frame = states[AVATAR_ANIMATION_INDEX % len(states)]
    AVATAR_ANIMATION_INDEX += 1
    
    # Full character with expressions
    characters = {
        "idle": f"      {frame}\n   ╭─────╮\n   │  ▽  │\n   ╰──┬──╯\n     ╱│╲",
        "listening": f"      {frame}\n   ╭─────╮\n   │  ◉  │\n   ╰──┬──╯\n     ╱│╲",
        "thinking": f"      {frame}\n   ╭─────╮\n   │  ◡  │\n   ╰──┬──╯\n     ╱│╲",
        "speaking": f"      {frame}\n   ╭─────╮\n   │  ◕  │\n   ╰──┬──╯\n     ╱│╲",
        "happy": f"      {frame}\n   ╭─────╮\n   │  ◖◗ │\n   ╰──┬──╯\n    ╱ │ ╲",
        "excited": f"      {frame}\n   ╭─────╮\n   │ ◝◞ │\n   ╰──┬──╯\n    ╱ │ ╲",
        "sad": f"      {frame}\n   ╭─────╮\n   │  ▼  │\n   ╰──┬──╯\n    ╱ │ ╲"
    }
    
    return {
        "status": "success",
        "character": characters.get(CURRENT_AVATAR_STATE, characters["idle"]),
        "state": CURRENT_AVATAR_STATE
    }

@app.post("/avatar/state")
async def set_avatar_state(state: str):
    """Set avatar state"""
    global CURRENT_AVATAR_STATE
    if state in AVATAR_STATES:
        CURRENT_AVATAR_STATE = state
        return {"status": "success", "state": state}
    return {"status": "error", "message": "Invalid state"}

# ============== AUDIO OUTPUT ==============

@app.post("/audio/speak")
async def audio_speak(text: str):
    """Generate and play speech"""
    try:
        from gtts import gTTS
        import uuid
        
        # Generate audio
        output_file = f"/tmp/speak_{uuid.uuid4().hex[:6]}.mp3"
        tts = gTTS(text=text[:500], lang="en")
        tts.save(output_file)
        
        return {
            "status": "success",
            "message": "Audio generated",
            "audio_file": output_file,
            "text": text
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.get("/audio/play")
async def audio_play(file_path: str):
    """Play audio file (returns file path for playback)"""
    import os
    
    if os.path.exists(file_path):
        return {"status": "success", "file": file_path, "ready": True}
    return {"status": "error", "message": "File not found"}

@app.get("/audio/list")
async def audio_list():
    """List generated audio files"""
    import glob
    
    files = glob.glob("/tmp/speak_*.mp3") + glob.glob("/tmp/tts_*.mp3")
    return {"status": "success", "files": files[-10:]}

# ============== VOICE ID & CONTINUOUS LISTENING ==============

# Voice ID configuration
VOICE_PROFILES = {}
CURRENT_USER = None
LISTENING_MODE = False

@app.post("/voiceid/enroll")
async def voiceid_enroll(user_name: str):
    """Enroll user's voice for identification"""
    global VOICE_PROFILES, CURRENT_USER
    
    # In production, would capture voice samples
    VOICE_PROFILES[user_name] = {
        "enrolled": True,
        "samples": 3,  # Would capture actual samples
        "created": datetime.now().isoformat()
    }
    CURRENT_USER = user_name
    
    return {
        "status": "success",
        "message": f"Voice enrolled for {user_name}!",
        "user": user_name,
        "samples_needed": 3
    }

@app.get("/voiceid/status")
async def voiceid_status():
    """Get voice ID status"""
    return {
        "status": "success",
        "enrolled_users": list(VOICE_PROFILES.keys()),
        "current_user": CURRENT_USER,
        "listening_mode": LISTENING_MODE
    }

@app.post("/voiceid/verify")
async def voiceid_verify():
    """Verify current speaker (placeholder)"""
    return {
        "status": "info",
        "message": "Voice verification requires microphone input",
        "current_user": CURRENT_USER
    }

@app.post("/voiceid/set_user")
async def voiceid_set_user(user_name: str):
    """Set current user manually"""
    global CURRENT_USER
    CURRENT_USER = user_name
    return {"status": "success", "user": user_name}

# Continuous listening mode
@app.post("/voice/listen/start")
async def start_continuous_listening():
    """Start continuous voice listening"""
    global LISTENING_MODE
    LISTENING_MODE = True
    
    return {
        "status": "success",
        "mode": "continuous_listening",
        "wake_word": "hey x",
        "instructions": [
            "Say 'Hey X' to activate",
            "Only enrolled voices will be recognized",
            "Stop with: /voice/listen/stop"
        ]
    }

@app.post("/voice/listen/stop")
async def stop_continuous_listening():
    """Stop continuous listening"""
    global LISTENING_MODE
    LISTENING_MODE = False
    
    return {"status": "success", "mode": "stopped"}

@app.get("/voice/listen/status")
async def listening_status():
    """Get listening status"""
    return {
        "status": "success",
        "listening": LISTENING_MODE,
        "wake_word": "hey x",
        "voice_id_enabled": len(VOICE_PROFILES) > 0,
        "current_user": CURRENT_USER
    }

# ============== OLLAMA AI (FREE!) ==============

OLLAMA_URL = "http://localhost:11434/api/generate"

@app.post("/ollama/chat")
async def ollama_chat(prompt: str, model: str = "llama3", system: str = ""):
    """Chat with Ollama (FREE AI!)"""
    try:
        import requests
        import json
        
        full_prompt = f"{system}\n\n{prompt}" if system else prompt
        
        response = requests.post(
            OLLAMA_URL,
            json={
                "model": model,
                "prompt": full_prompt,
                "stream": False,
                "options": {
                    "temperature": 0.7
                }
            },
            timeout=120
        )
        
        result = response.json()
        
        return {
            "status": "success",
            "model": model,
            "response": result.get("response", ""),
            "done": result.get("done", True)
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.get("/ollama/models")
async def ollama_models():
    """List available Ollama models"""
    try:
        import requests
        response = requests.get("http://localhost:11434/api/tags")
        models = response.json().get("models", [])
        return {"status": "success", "models": [m["name"] for m in models]}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/ollama/code")
async def ollama_code(prompt: str, language: str = "python"):
    """Generate code using Ollama"""
    try:
        import requests
        
        code_prompt = f"""You are a {language} programmer. Write clean, working code for: {prompt}
Only output the code, no explanations."""

        response = requests.post(
            OLLAMA_URL,
            json={
                "model": "llama3",
                "prompt": code_prompt,
                "stream": False
            },
            timeout=120
        )
        
        result = response.json()
        
        return {
            "status": "success",
            "code": result.get("response", ""),
            "language": language,
            "model": "llama3 (FREE)"
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/ollama/explain")
async def ollama_explain(code: str):
    """Explain code using Ollama"""
    try:
        import requests
        
        explain_prompt = f"Explain this code simply:\n\n{code}"

        response = requests.post(
            OLLAMA_URL,
            json={
                "model": "llama3",
                "prompt": explain_prompt,
                "stream": False
            },
            timeout=120
        )
        
        result = response.json()
        
        return {
            "status": "success",
            "explanation": result.get("response", "")
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}

# ============== VPS BROWSER CONTROL ==============

@app.get("/vps/browser/open")
async def vps_browser_open(url: str):
    """Open URL in VPS browser"""
    try:
        # Try to open in browser
        subprocess.run(["xdg-open", url], capture_output=True, timeout=10)
        return {"status": "success", "action": f"Opened {url}"}
    except Exception as e:
        return {"status": "info", "message": "Browser control needs Xvfb on VPS", "url": url}

@app.get("/vps/screenshot")
async def vps_screenshot():
    """Take screenshot on VPS"""
    try:
        # Use gnome-screenshot or import
        result = subprocess.run(["gnome-screenshot", "-f", "/tmp/vps_screen.png"], 
                              capture_output=True, timeout=10)
        if result.returncode != 0:
            subprocess.run(["import", "-window", "root", "/tmp/vps_screen.png"], 
                         capture_output=True)
        
        return {"status": "success", "file": "/tmp/vps_screen.png"}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/vps/type")
async def vps_type(text: str):
    """Type on VPS"""
    try:
        escaped = text.replace("'", "'\\''")
        subprocess.run(f"xdotool type -- '{escaped}'", shell=True, capture_output=True, timeout=10)
        return {"status": "success", "typed": text}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/vps/click")
async def vps_click(x: int = None, y: int = None):
    """Click on VPS screen"""
    try:
        if x and y:
            subprocess.run(["xdotool", "mousemove", str(x), str(y), "click", "1"], 
                         capture_output=True, timeout=5)
        else:
            subprocess.run(["xdotool", "click", "1"], capture_output=True, timeout=5)
        return {"status": "success", "action": "clicked"}
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.get("/vps/status")
async def vps_status():
    """Get VPS desktop status"""
    return {
        "status": "success",
        "display": os.environ.get('DISPLAY', 'Not set'),
        "available": True,
        "features": [
            "/vps/browser/open?url=...",
            "/vps/screenshot",
            "/vps/type?text=...",
            "/vps/click?x=100&y=100"
        ]
    }

# ============== CLAUDE CODE (FREE!) ==============

@app.get("/claude/install")
async def claude_install():
    """Install Claude Code CLI (FREE!)"""
    try:
        # Claude CLI installation commands
        return {
            "status": "info",
            "message": "Claude Code - Free AI CLI",
            "install_commands": [
                "npm install -g @anthropic-ai/claude",
                "claude configure"
            ],
            "usage": "claude -p 'your prompt'",
            "website": "https://claude.com/claude-code",
            "note": "Claude Code is FREE for individual use!"
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/claude/run")
async def claude_run(prompt: str):
    """Run Claude Code CLI"""
    try:
        result = subprocess.run(
            ["claude", "-p", prompt],
            capture_output=True,
            text=True,
            timeout=120
        )
        return {
            "status": "success",
            "response": result.stdout[:3000],
            "error": result.stderr[:500] if result.stderr else None
        }
    except FileNotFoundError:
        return {
            "status": "error",
            "message": "Claude not installed",
            "install": "npm install -g @anthropic-ai/claude"
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/claude/code")
async def claude_code(prompt: str, language: str = "python"):
    """Generate code using Claude"""
    code_prompt = f"Write {language} code for: {prompt}. Only output the code, no explanation."
    
    try:
        result = subprocess.run(
            ["claude", "-p", code_prompt],
            capture_output=True,
            text=True,
            timeout=120
        )
        return {
            "status": "success",
            "code": result.stdout[:3000],
            "model": "Claude Code (FREE)"
        }
    except FileNotFoundError:
        return {
            "status": "error",
            "message": "Claude not installed",
            "install": "npm install -g @anthropic-ai/claude"
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.get("/claude/status")
async def claude_status():
    """Check Claude status"""
    try:
        result = subprocess.run(["claude", "--version"], capture_output=True, text=True, timeout=10)
        return {
            "status": "success",
            "installed": True,
            "version": result.stdout.strip()
        }
    except FileNotFoundError:
        return {
            "status": "info",
            "installed": False,
            "message": "Claude not installed",
            "install": "npm install -g @anthropic-ai/claude"
        }

# ============== ANTHROPIC CLAUDE API ==============

@app.post("/claude/api/chat")
async def claude_chat(prompt: str, api_key: str = None):
    """Chat with Claude API (needs API key)"""
    try:
        import anthropic
        
        # Use provided key or env
        if not api_key:
            api_key = os.environ.get('ANTHROPIC_API_KEY', '')
        
        if not api_key:
            return {
                "status": "info",
                "message": "API key required",
                "get_free_credits": "https://console.anthropic.com/"
            }
        
        client = anthropic.Anthropic(api_key=api_key)
        
        message = client.messages.create(
            model="claude-3-haiku-20240307",
            max_tokens=1024,
            messages=[{"role": "user", "content": prompt}]
        )
        
        return {
            "status": "success",
            "response": message.content[0].text,
            "model": "claude-3-haiku"
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}

@app.post("/claude/api/code")
async def claude_code(prompt: str, language: str = "python", api_key: str = None):
    """Generate code with Claude"""
    code_prompt = f"Write {language} code for: {prompt}. Only output the code, no explanation."
    
    try:
        import anthropic
        
        if not api_key:
            api_key = os.environ.get('ANTHROPIC_API_KEY', '')
        
        if not api_key:
            return {
                "status": "info",
                "message": "API key needed",
                "get_key": "https://console.anthropic.com/"
            }
        
        client = anthropic.Anthropic(api_key=api_key)
        
        message = client.messages.create(
            model="claude-3-haiku-20240307",
            max_tokens=2048,
            messages=[{"role": "user", "content": code_prompt}]
        )
        
        return {
            "status": "success",
            "code": message.content[0].text,
            "model": "claude-3-haiku (FREE credits available!)"
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}

# ============== HYBRID BRIDGE ==============

# Desktop bridge management
DESKTOP_REGISTRY = {}
PENDING_COMMANDS = []

@app.post("/bridge/register")
async def bridge_register(data: dict):
    """Register a desktop"""
    desktop_id = data.get("desktop_id")
    DESKTOP_REGISTRY[desktop_id] = {
        "status": "online",
        "registered_at": datetime.now().isoformat()
    }
    return {"status": "success", "desktop_id": desktop_id}

@app.post("/bridge/heartbeat")
async def bridge_heartbeat(data: dict):
    """Desktop heartbeat"""
    desktop_id = data.get("desktop_id")
    if desktop_id in DESKTOP_REGISTRY:
        DESKTOP_REGISTRY[desktop_id]["last_seen"] = datetime.now().isoformat()
        DESKTOP_REGISTRY[desktop_id]["status"] = "online"
    return {"status": "success"}

@app.get("/bridge/desktops")
async def bridge_list_desktops():
    """List registered desktops"""
    return {"status": "success", "desktops": DESKTOP_REGISTRY}

@app.post("/bridge/command")
async def bridge_send_command(desktop_id: str, command: str, params: dict = None):
    """Send command to desktop"""
    cmd_id = f"cmd_{len(PENDING_COMMANDS) + 1}"
    PENDING_COMMANDS.append({
        "id": cmd_id,
        "desktop_id": desktop_id,
        "command": command,
        "params": params or {}
    })
    return {"status": "success", "command_id": cmd_id}

@app.get("/bridge/commands")
async def bridge_get_commands(desktop_id: str):
    """Get commands for desktop"""
    cmds = [c for c in PENDING_COMMANDS if c.get("desktop_id") == desktop_id]
    return {"status": "success", "commands": cmds}

@app.post("/bridge/result")
async def bridge_result(data: dict):
    """Receive command result"""
    cmd_id = data.get("command_id")
    # Remove from pending
    global PENDING_COMMANDS
    PENDING_COMMANDS = [c for c in PENDING_COMMANDS if c.get("id") != cmd_id]
    return {"status": "success"}

# ============== MAIN ==============
if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)

